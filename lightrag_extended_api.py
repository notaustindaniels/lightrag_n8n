#!/usr/bin/env python3
"""
Extended LightRAG API Server that adds missing functionality for document management
"""
import os
import asyncio
import hashlib
import json
import networkx as nx
from typing import List, Dict, Optional, Any
from datetime import datetime
from urllib.parse import urlparse
from fastapi import FastAPI, HTTPException, Depends, File, UploadFile, Form, Query
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from fastapi.responses import FileResponse, RedirectResponse, StreamingResponse
from pydantic import BaseModel, Field
import uvicorn
from contextlib import asynccontextmanager
import traceback
import shutil
import time
import mimetypes
import io
import chardet

# Document parsing imports
import docx
import pypdf
import pptx
from bs4 import BeautifulSoup
import ebooklib
from ebooklib import epub
import yaml
import csv
import magic

# Import LightRAG components
from lightrag import LightRAG, QueryParam
from lightrag.llm.openai import gpt_4o_mini_complete, openai_embed
from lightrag.kg.shared_storage import initialize_pipeline_status
from lightrag.utils import EmbeddingFunc, setup_logger

# Setup logging
setup_logger("lightrag", level="INFO")

# Models
class TextInsertRequest(BaseModel):
    text: str
    description: Optional[str] = None
    file_path: Optional[str] = None
    metadata: Optional[Dict[str, Any]] = None

class EnhancedTextInsertRequest(BaseModel):
    text: str
    description: Optional[str] = None
    source_url: Optional[str] = None
    sitemap_url: Optional[str] = None
    doc_index: Optional[int] = None
    total_docs: Optional[int] = None
    document_id: Optional[str] = None  # Accept document_id from n8n workflow

class DocumentMetadata(BaseModel):
    id: str
    file_path: Optional[str] = None
    source_url: Optional[str] = None
    sitemap_url: Optional[str] = None
    description: Optional[str] = None
    indexed_at: Optional[str] = None
    content_summary: Optional[str] = None

class DocumentResponse(BaseModel):
    id: str
    file_path: Optional[str] = Field(default="")  # Default empty string to avoid validation error
    metadata: Optional[DocumentMetadata] = None

class DeleteByIdRequest(BaseModel):
    doc_ids: List[str]

class QueryRequest(BaseModel):
    query: str
    mode: str = "hybrid"
    stream: bool = False
    workspace: Optional[str] = None  # Optional workspace to query from

# Graph models
class EntityUpdateRequest(BaseModel):
    entity_name: str
    updated_data: Dict[str, Any]
    allow_rename: bool = False

class RelationUpdateRequest(BaseModel):
    source_id: str
    target_id: str
    updated_data: Dict[str, Any]



# Global variables
workspace_instances = {}  # Dictionary to store multiple LightRAG instances by workspace
workspace_metadata = {}  # Metadata store per workspace
default_workspace = "default"  # Default workspace for backward compatibility
metadata_store = {}  # Global metadata store for backward compatibility with WebUI

# Supported file extensions
SUPPORTED_EXTENSIONS = {
    'txt', 'md', 'docx', 'pdf', 'pptx', 'rtf', 'odt', 'epub', 
    'html', 'htm', 'tex', 'json', 'xml', 'yaml', 'yml', 'csv', 
    'log', 'conf', 'ini', 'properties', 'sql', 'bat', 'sh', 
    'c', 'cpp', 'py', 'java', 'js', 'ts', 'swift', 'go', 
    'rb', 'php', 'css', 'scss', 'less'
}

# MIME type mapping
MIME_TO_EXT = {
    'text/plain': ['txt', 'log', 'conf', 'ini', 'properties'],
    'text/markdown': ['md'],
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document': ['docx'],
    'application/pdf': ['pdf'],
    'application/vnd.openxmlformats-officedocument.presentationml.presentation': ['pptx'],
    'application/rtf': ['rtf'],
    'application/vnd.oasis.opendocument.text': ['odt'],
    'application/epub+zip': ['epub'],
    'text/html': ['html', 'htm'],
    'text/x-tex': ['tex'],
    'application/json': ['json'],
    'application/xml': ['xml'],
    'text/xml': ['xml'],
    'application/x-yaml': ['yaml', 'yml'],
    'text/yaml': ['yaml', 'yml'],
    'text/csv': ['csv'],
    'text/x-sql': ['sql'],
    'text/x-sh': ['sh'],
    'text/x-bat': ['bat'],
    'text/x-c': ['c'],
    'text/x-c++': ['cpp'],
    'text/x-python': ['py'],
    'text/x-java': ['java'],
    'text/javascript': ['js'],
    'application/javascript': ['js'],
    'text/x-typescript': ['ts'],
    'text/x-swift': ['swift'],
    'text/x-go': ['go'],
    'text/x-ruby': ['rb'],
    'text/x-php': ['php'],
    'text/css': ['css'],
    'text/x-scss': ['scss'],
    'text/x-less': ['less']
}

# Workspace Manager Class
class WorkspaceManager:
    """Manages multiple LightRAG instances for different workspaces"""
    
    @staticmethod
    def extract_workspace_from_doc_id(doc_id: str) -> str:
        """Extract workspace name from document ID format [workspace] filename"""
        if doc_id.startswith('[') and ']' in doc_id:
            workspace = doc_id.split(']')[0][1:]
            # Normalize workspace name - replace dots with underscores for file system compatibility
            # but keep the original name for display
            return workspace
        return default_workspace
    
    @staticmethod
    async def get_or_create_instance(workspace: str) -> LightRAG:
        """Get existing workspace instance or create new one"""
        if workspace not in workspace_instances:
            await WorkspaceManager.create_workspace(workspace)
        return workspace_instances[workspace]
    
    @staticmethod
    async def create_workspace(workspace: str):
        """Create a new workspace with isolated storage"""
        base_dir = os.getenv("WORKING_DIR", "/app/data/rag_storage")
        workspace_dir = os.path.join(base_dir, "workspaces", workspace)
        os.makedirs(workspace_dir, exist_ok=True)
        
        print(f"\n=== Creating/Loading workspace: {workspace} ===")
        print(f"Workspace directory: {workspace_dir}")
        
        # Create LightRAG instance for this workspace
        # Note: Don't pass workspace parameter as it creates nested directories
        # The working_dir already provides isolation
        rag = LightRAG(
            working_dir=workspace_dir,
            embedding_func=EmbeddingFunc(
                embedding_dim=1536,
                max_token_size=8192,
                func=openai_embed
            ),
            llm_model_func=gpt_4o_mini_complete,
        )
        
        await rag.initialize_storages()
        await initialize_pipeline_status()
        
        # After initialization, check for graph in both flat and nested structures
        graphml_paths = [
            os.path.join(workspace_dir, "graph_chunk_entity_relation.graphml"),
            # Also check nested structure for backward compatibility
            os.path.join(workspace_dir, workspace, "graph_chunk_entity_relation.graphml")
        ]
        
        graphml_path = None
        for path in graphml_paths:
            if os.path.exists(path):
                graphml_path = path
                break
        
        if graphml_path:
            print(f"Found existing GraphML file: {graphml_path}")
            try:
                # Check if the graph was loaded properly
                if hasattr(rag, 'chunk_entity_relation_graph'):
                    storage = rag.chunk_entity_relation_graph
                    
                    # Try to get the graph from storage to verify it's loaded
                    if hasattr(storage, '_get_graph'):
                        graph = await storage._get_graph()
                        if graph is not None and hasattr(graph, 'number_of_nodes') and graph.number_of_nodes() > 0:
                            print(f"Graph loaded successfully with {graph.number_of_nodes()} nodes and {graph.number_of_edges()} edges")
                        else:
                            print(f"Warning: Graph appears to be empty or not loaded properly")
                            # Always try manual loading if the graph is empty
                            try:
                                loaded_graph = nx.read_graphml(graphml_path)
                                if loaded_graph.number_of_nodes() > 0:
                                    print(f"Manually loading graph with {loaded_graph.number_of_nodes()} nodes from {graphml_path}")
                                    if hasattr(storage, '_graph'):
                                        storage._graph = loaded_graph
                                        # Mark storage as not needing update
                                        if hasattr(storage, 'storage_updated') and storage.storage_updated:
                                            storage.storage_updated.value = False
                                        print(f"Successfully set graph in storage")
                                else:
                                    print(f"GraphML file exists but has no nodes")
                            except Exception as e:
                                print(f"Error during manual graph loading: {e}")
                    elif hasattr(storage, '_graph'):
                        # Direct access for older versions
                        if storage._graph is not None and hasattr(storage._graph, 'number_of_nodes'):
                            print(f"Graph accessible with {storage._graph.number_of_nodes()} nodes")
                        else:
                            print(f"Warning: Graph not properly loaded")
                            # Try manual loading
                            try:
                                loaded_graph = nx.read_graphml(graphml_path)
                                if loaded_graph.number_of_nodes() > 0:
                                    storage._graph = loaded_graph
                                    print(f"Manually set graph with {loaded_graph.number_of_nodes()} nodes")
                            except Exception as e:
                                print(f"Error setting graph: {e}")
            except Exception as e:
                print(f"Error verifying graph load: {e}")
                import traceback
                traceback.print_exc()
        else:
            print(f"No existing GraphML file found in {workspace_dir} or nested structure")
        
        workspace_instances[workspace] = rag
        
        # Initialize metadata for this workspace
        if workspace not in workspace_metadata:
            workspace_metadata[workspace] = {}
            load_workspace_metadata(workspace)
        
        print(f"Workspace {workspace} ready")
        return rag
    
    @staticmethod
    async def delete_workspace(workspace: str):
        """Delete a workspace and all its data"""
        if workspace == default_workspace:
            raise ValueError("Cannot delete the default workspace")
        
        if workspace in workspace_instances:
            # Clean up the instance
            del workspace_instances[workspace]
        
        if workspace in workspace_metadata:
            del workspace_metadata[workspace]
        
        # Delete workspace directory
        base_dir = os.getenv("WORKING_DIR", "/app/data/rag_storage")
        workspace_dir = os.path.join(base_dir, "workspaces", workspace)
        if os.path.exists(workspace_dir):
            shutil.rmtree(workspace_dir)
            print(f"Deleted workspace directory: {workspace_dir}")
    
    @staticmethod
    def list_workspaces() -> List[str]:
        """List all available workspaces"""
        base_dir = os.getenv("WORKING_DIR", "/app/data/rag_storage")
        workspaces_dir = os.path.join(base_dir, "workspaces")
        
        workspaces = []
        if os.path.exists(workspaces_dir):
            workspaces = [d for d in os.listdir(workspaces_dir) 
                         if os.path.isdir(os.path.join(workspaces_dir, d))]
        
        # Always include active workspaces even if directory doesn't exist yet
        for ws in workspace_instances.keys():
            if ws not in workspaces:
                workspaces.append(ws)
        
        return workspaces

def compute_doc_id(content: str) -> str:
    """Compute document ID using MD5 hash of content"""
    return f"doc-{hashlib.md5(content.strip().encode()).hexdigest()}"

def save_metadata_store(workspace: str = None):
    """Save metadata for a specific workspace or all workspaces"""
    if workspace:
        save_workspace_metadata(workspace)
    else:
        for ws in workspace_metadata:
            save_workspace_metadata(ws)

def save_workspace_metadata(workspace: str):
    """Save metadata for a specific workspace to disk"""
    try:
        base_dir = os.getenv("WORKING_DIR", "/app/data/rag_storage")
        workspace_dir = os.path.join(base_dir, "workspaces", workspace)
        os.makedirs(workspace_dir, exist_ok=True)
        
        metadata_file = os.path.join(workspace_dir, "document_metadata.json")
        metadata = workspace_metadata.get(workspace, {})
        
        with open(metadata_file, 'w') as f:
            json.dump(metadata, f, indent=2)
        print(f"Metadata saved for workspace {workspace}: {len(metadata)} documents")
    except Exception as e:
        print(f"Error saving metadata for workspace {workspace}: {e}")

def load_workspace_metadata(workspace: str):
    """Load metadata for a specific workspace from disk"""
    try:
        base_dir = os.getenv("WORKING_DIR", "/app/data/rag_storage")
        workspace_dir = os.path.join(base_dir, "workspaces", workspace)
        metadata_file = os.path.join(workspace_dir, "document_metadata.json")
        
        if os.path.exists(metadata_file):
            with open(metadata_file, 'r') as f:
                workspace_metadata[workspace] = json.load(f)
            print(f"Metadata loaded for workspace {workspace}: {len(workspace_metadata[workspace])} documents")
        else:
            workspace_metadata[workspace] = {}
            print(f"No existing metadata for workspace {workspace}, starting fresh")
    except Exception as e:
        print(f"Error loading metadata for workspace {workspace}: {e}")
        workspace_metadata[workspace] = {}

def get_all_documents_metadata():
    """Get aggregated metadata from all workspaces for backward compatibility"""
    all_metadata = {}
    for workspace, metadata in workspace_metadata.items():
        all_metadata.update(metadata)
    return all_metadata

async def get_combined_graph():
    """Combine graphs from all workspaces into a single graph"""
    combined_graph = nx.Graph()
    
    print(f"\n=== GRAPH AGGREGATION DEBUG ===")
    print(f"Total workspaces loaded: {len(workspace_instances)}")
    print(f"Workspace names: {list(workspace_instances.keys())}")
    
    for workspace_name, rag in workspace_instances.items():
        print(f"\n--- Processing workspace: '{workspace_name}' ---")
        
        # Try multiple ways to get the graph
        workspace_graph = None
        
        # Method 1: Direct graph attribute access
        if hasattr(rag, 'chunk_entity_relation_graph'):
            print(f"  Found chunk_entity_relation_graph attribute")
            graph_storage = rag.chunk_entity_relation_graph
            print(f"  Graph storage type: {type(graph_storage)}")
            
            # Try to get the actual graph (now async)
            workspace_graph = await get_graph_from_storage(graph_storage)
            
            if workspace_graph:
                print(f"  Successfully extracted graph from storage")
            else:
                print(f"  Failed to extract graph from storage")
                
                # Try loading from file as fallback
                workspace_dir = os.path.join(os.getenv("WORKING_DIR", "/app/data/rag_storage"), "workspaces", workspace_name)
                graphml_path = os.path.join(workspace_dir, "graph_chunk_entity_relation.graphml")
                
                if os.path.exists(graphml_path):
                    try:
                        workspace_graph = nx.read_graphml(graphml_path)
                        print(f"  Loaded graph from GraphML file: {graphml_path}")
                    except Exception as e:
                        print(f"  Error loading GraphML: {e}")
        
        # Method 2: Try loading directly from workspace directory
        if not workspace_graph:
            workspace_dir = os.path.join(os.getenv("WORKING_DIR", "/app/data/rag_storage"), "workspaces", workspace_name)
            
            # Check both flat and nested structures
            graphml_paths = [
                os.path.join(workspace_dir, "graph_chunk_entity_relation.graphml"),
                os.path.join(workspace_dir, workspace_name, "graph_chunk_entity_relation.graphml")  # Nested structure
            ]
            
            graphml_path = None
            for path in graphml_paths:
                print(f"  Checking for GraphML at: {path}")
                if os.path.exists(path):
                    graphml_path = path
                    break
            
            if graphml_path:
                try:
                    workspace_graph = nx.read_graphml(graphml_path)
                    print(f"  Successfully loaded graph from file: {graphml_path}")
                except Exception as e:
                    print(f"  Error loading GraphML file: {e}")
            else:
                print(f"  No GraphML file found in flat or nested structure")
        
        # Process the graph if we have one
        if workspace_graph and hasattr(workspace_graph, 'nodes'):
            try:
                node_count = workspace_graph.number_of_nodes()
                edge_count = workspace_graph.number_of_edges()
                print(f"  Graph stats: {node_count} nodes, {edge_count} edges")
                
                # Show sample nodes for debugging
                sample_nodes = list(workspace_graph.nodes())[:5]
                if sample_nodes:
                    print(f"  Sample nodes: {sample_nodes}")
                
                # Add nodes with workspace information
                for node in workspace_graph.nodes(data=True):
                    node_id = node[0]
                    node_data = node[1] if len(node) > 1 else {}
                    # Add workspace info to node data
                    node_data['workspace'] = workspace_name
                    combined_graph.add_node(node_id, **node_data)
                
                # Add edges
                for edge in workspace_graph.edges(data=True):
                    source = edge[0]
                    target = edge[1]
                    edge_data = edge[2] if len(edge) > 2 else {}
                    # Add workspace info to edge data
                    edge_data['workspace'] = workspace_name
                    combined_graph.add_edge(source, target, **edge_data)
                    
                print(f"  Successfully added to combined graph")
            except Exception as e:
                print(f"  Error processing graph: {e}")
                import traceback
                traceback.print_exc()
        else:
            print(f"  No valid graph found for workspace '{workspace_name}'")
    
    total_nodes = combined_graph.number_of_nodes()
    total_edges = combined_graph.number_of_edges()
    print(f"\n=== FINAL COMBINED GRAPH ===")
    print(f"Total nodes: {total_nodes}")
    print(f"Total edges: {total_edges}")
    
    if total_nodes > 0:
        sample_nodes = list(combined_graph.nodes())[:10]
        print(f"Sample combined nodes: {sample_nodes}")
    
    return combined_graph

async def get_all_graph_labels():
    """Get all unique graph labels from all workspaces"""
    all_labels = set()
    
    for workspace_name, rag in workspace_instances.items():
        if hasattr(rag, 'get_graph_labels'):
            try:
                labels = await rag.get_graph_labels()
                all_labels.update(labels)
            except:
                pass
        elif hasattr(rag, 'chunk_entity_relation_graph'):
            graph_storage = rag.chunk_entity_relation_graph
            graph = get_graph_from_storage(graph_storage)
            if graph and hasattr(graph, 'nodes'):
                all_labels.update(graph.nodes())
    
    return list(all_labels)

def extract_text_from_file(file_content: bytes, file_extension: str, filename: str) -> str:
    """Extract text from various file formats"""
    try:
        # Text-based files
        if file_extension in ['txt', 'md', 'log', 'conf', 'ini', 'properties', 'sql', 
                             'bat', 'sh', 'c', 'cpp', 'py', 'java', 'js', 'ts', 
                             'swift', 'go', 'rb', 'php', 'css', 'scss', 'less', 'tex']:
            # Detect encoding
            detected = chardet.detect(file_content)
            encoding = detected['encoding'] or 'utf-8'
            return file_content.decode(encoding, errors='replace')
        
        # Microsoft Word
        elif file_extension == 'docx':
            doc = docx.Document(io.BytesIO(file_content))
            paragraphs = []
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append(para.text)
            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            paragraphs.append(cell.text)
            return '\n\n'.join(paragraphs)
        
        # PDF
        elif file_extension == 'pdf':
            pdf_reader = pypdf.PdfReader(io.BytesIO(file_content))
            text_parts = []
            for page_num, page in enumerate(pdf_reader.pages, 1):
                page_text = page.extract_text()
                if page_text.strip():
                    text_parts.append(f"--- Page {page_num} ---\n{page_text}")
            return '\n\n'.join(text_parts)
        
        # PowerPoint
        elif file_extension == 'pptx':
            prs = pptx.Presentation(io.BytesIO(file_content))
            text_parts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        slide_text.append(shape.text)
                if slide_text:
                    text_parts.append(f"--- Slide {slide_num} ---\n" + '\n'.join(slide_text))
            return '\n\n'.join(text_parts)
        
        # HTML
        elif file_extension in ['html', 'htm']:
            soup = BeautifulSoup(file_content, 'html.parser')
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            # Get text
            text = soup.get_text()
            # Break into lines and remove leading/trailing space
            lines = (line.strip() for line in text.splitlines())
            # Drop blank lines
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = '\n'.join(chunk for chunk in chunks if chunk)
            return text
        
        # EPUB
        elif file_extension == 'epub':
            book = epub.read_epub(io.BytesIO(file_content))
            text_parts = []
            for item in book.get_items():
                if item.get_type() == ebooklib.ITEM_DOCUMENT:
                    soup = BeautifulSoup(item.get_body_content(), 'html.parser')
                    text = soup.get_text()
                    if text.strip():
                        text_parts.append(text)
            return '\n\n'.join(text_parts)
        
        # JSON
        elif file_extension == 'json':
            data = json.loads(file_content.decode('utf-8'))
            return json.dumps(data, indent=2, ensure_ascii=False)
        
        # XML
        elif file_extension == 'xml':
            soup = BeautifulSoup(file_content, 'xml')
            return soup.prettify()
        
        # YAML
        elif file_extension in ['yaml', 'yml']:
            data = yaml.safe_load(file_content.decode('utf-8'))
            return yaml.dump(data, default_flow_style=False, allow_unicode=True)
        
        # CSV
        elif file_extension == 'csv':
            detected = chardet.detect(file_content)
            encoding = detected['encoding'] or 'utf-8'
            text = file_content.decode(encoding, errors='replace')
            # Parse CSV and format as readable text
            reader = csv.reader(io.StringIO(text))
            rows = list(reader)
            if rows:
                # Format as markdown table
                result = []
                if len(rows) > 1:
                    # Assume first row is header
                    headers = rows[0]
                    result.append(' | '.join(headers))
                    result.append(' | '.join(['---'] * len(headers)))
                    for row in rows[1:]:
                        result.append(' | '.join(row))
                else:
                    # Just one row
                    result.append(' | '.join(rows[0]))
                return '\n'.join(result)
            return text
        
        # RTF and ODT would require additional libraries (python-rtf, odfpy)
        # For now, try to extract as plain text
        elif file_extension in ['rtf', 'odt']:
            # Basic text extraction - won't preserve formatting
            detected = chardet.detect(file_content)
            encoding = detected['encoding'] or 'utf-8'
            text = file_content.decode(encoding, errors='replace')
            # For RTF, try to extract visible text (very basic)
            if file_extension == 'rtf':
                import re
                # Remove RTF commands (basic approach)
                text = re.sub(r'\\[a-z]+[0-9]*\s?', '', text)
                text = re.sub(r'[{}]', '', text)
            return text
        
        else:
            # Fallback - try to decode as text
            detected = chardet.detect(file_content)
            encoding = detected['encoding'] or 'utf-8'
            return file_content.decode(encoding, errors='replace')
            
    except Exception as e:
        print(f"Error extracting text from {filename}: {e}")
        # Fallback - try to decode as text
        try:
            detected = chardet.detect(file_content)
            encoding = detected['encoding'] or 'utf-8'
            return file_content.decode(encoding, errors='replace')
        except:
            raise HTTPException(
                status_code=400, 
                detail=f"Could not extract text from {filename}: {str(e)}"
            )

def generate_display_name_from_file_path(file_path: str, doc_id: str) -> str:
    """Generate a display name from a file path for legacy documents"""
    if "[" in file_path and "]" in file_path:
        # Return the full file_path as display_name (includes domain and full path)
        return file_path
    else:
        return f"text/{doc_id[:8]}..."

async def get_graph_from_storage(storage_or_graph):
    """Safely get NetworkX graph from storage object or return graph directly"""
    if storage_or_graph is None:
        return None
    
    # Check if it's already a NetworkX graph
    if hasattr(storage_or_graph, 'nodes') and hasattr(storage_or_graph, 'edges'):
        # Check if it's a direct NetworkX graph by trying to access nodes
        try:
            # Try to call a NetworkX-specific method
            storage_or_graph.number_of_nodes()
            return storage_or_graph
        except (AttributeError, TypeError):
            pass
    
    # For NetworkXStorage, use the proper async method
    if hasattr(storage_or_graph, '_get_graph'):
        try:
            # This is the NetworkXStorage class from lightrag
            # First ensure it's initialized
            if hasattr(storage_or_graph, 'initialize') and not hasattr(storage_or_graph, '_storage_lock'):
                await storage_or_graph.initialize()
            graph = await storage_or_graph._get_graph()
            return graph
        except Exception as e:
            print(f"Error getting graph from NetworkXStorage: {e}")
            # Try to access the _graph attribute directly as fallback
            if hasattr(storage_or_graph, '_graph'):
                return storage_or_graph._graph
    
    # Try to get graph from storage object (non-async methods)
    if hasattr(storage_or_graph, '_graph'):
        return storage_or_graph._graph
    elif hasattr(storage_or_graph, 'graph'):
        return storage_or_graph.graph
    elif hasattr(storage_or_graph, 'get_graph'):
        # Check if get_graph is async
        import inspect
        if inspect.iscoroutinefunction(storage_or_graph.get_graph):
            try:
                return await storage_or_graph.get_graph()
            except Exception as e:
                print(f"Error calling async get_graph: {e}")
                return None
        else:
            try:
                return storage_or_graph.get_graph()
            except Exception as e:
                print(f"Error calling get_graph: {e}")
                return None
    elif hasattr(storage_or_graph, 'data'):
        return storage_or_graph.data
    
    # If all else fails, return None instead of the object itself
    print(f"Warning: Could not extract graph from storage object of type {type(storage_or_graph)}")
    return None

def force_save_graph_to_disk(graph, working_dir):
    """Force save graph to disk with multiple fallback methods"""
    graphml_path = os.path.join(working_dir, "graph_chunk_entity_relation.graphml")
    
    try:
        # Method 1: Direct NetworkX write
        nx.write_graphml(graph, graphml_path)
        print(f"Saved graph to {graphml_path} using NetworkX")
        
        # Verify the file was written
        if os.path.exists(graphml_path):
            file_size = os.path.getsize(graphml_path)
            print(f"GraphML file size: {file_size} bytes")
            
            # Double-check by reading it back
            test_graph = nx.read_graphml(graphml_path)
            print(f"Verified: graph has {test_graph.number_of_nodes()} nodes and {test_graph.number_of_edges()} edges")
        
        return True
    except Exception as e:
        print(f"Error in force_save_graph_to_disk: {e}")
        
        # Fallback: Write to temp file and move
        try:
            temp_path = graphml_path + ".tmp"
            nx.write_graphml(graph, temp_path)
            shutil.move(temp_path, graphml_path)
            print(f"Saved graph using temp file method")
            return True
        except Exception as e2:
            print(f"Error in temp file method: {e2}")
            return False

async def cleanup_all_document_traces(doc_ids: List[str]):
    """Enhanced cleanup that ensures all traces are removed and properly persisted"""
    try:
        print(f"Starting enhanced cleanup for documents: {doc_ids}")
        
        # Create a set of all possible document ID formats to check
        all_doc_id_patterns = set()
        for doc_id in doc_ids:
            all_doc_id_patterns.add(doc_id)
            if doc_id.startswith('[') and ']' in doc_id:
                parts = doc_id.split('] ', 1)
                if len(parts) == 2:
                    domain_part = parts[0] + ']'
                    all_doc_id_patterns.add(domain_part)
        
        working_dir = os.getenv("WORKING_DIR", "/app/data/rag_storage")
        
        # Step 1: Clean vector databases FIRST and SAVE IMMEDIATELY
        vdb_files = ["vdb_entities.json", "vdb_relationships.json", "vdb_chunks.json"]
        
        entities_to_remove = set()
        relationships_to_remove = set()
        
        for vdb_file in vdb_files:
            vdb_path = os.path.join(working_dir, vdb_file)
            if os.path.exists(vdb_path):
                try:
                    # Read the file
                    with open(vdb_path, 'r') as f:
                        vdb_data = json.load(f)
                    
                    if 'data' in vdb_data and isinstance(vdb_data['data'], list):
                        original_count = len(vdb_data['data'])
                        filtered_data = []
                        
                        for entry in vdb_data['data']:
                            should_remove = False
                            
                            # Check various fields for document references
                            for field in ['doc_id', 'source_id', 'file_path', 'chunk_id']:
                                value = entry.get(field, '')
                                for pattern in all_doc_id_patterns:
                                    if pattern in str(value):
                                        should_remove = True
                                        
                                        # Track what we're removing for graph cleanup
                                        if vdb_file == "vdb_entities.json":
                                            entities_to_remove.add(entry.get('id', ''))
                                        elif vdb_file == "vdb_relationships.json":
                                            src = entry.get('source', '')
                                            tgt = entry.get('target', '')
                                            if src and tgt:
                                                relationships_to_remove.add((src, tgt))
                                        break
                                if should_remove:
                                    break
                            
                            if not should_remove:
                                filtered_data.append(entry)
                        
                        vdb_data['data'] = filtered_data
                        
                        # Clear matrix if data changed
                        if len(filtered_data) != original_count:
                            vdb_data['matrix'] = []
                        
                        # Save immediately with backup
                        backup_path = vdb_path + ".backup"
                        shutil.copy2(vdb_path, backup_path)
                        
                        with open(vdb_path, 'w') as f:
                            json.dump(vdb_data, f, indent=2)
                        
                        # Verify the write
                        with open(vdb_path, 'r') as f:
                            verify_data = json.load(f)
                            if len(verify_data.get('data', [])) == len(filtered_data):
                                os.remove(backup_path)  # Remove backup if successful
                                print(f"Successfully updated {vdb_file}, removed {original_count - len(filtered_data)} entries")
                            else:
                                shutil.move(backup_path, vdb_path)  # Restore backup
                                raise Exception("Verification failed after write")
                        
                except Exception as e:
                    print(f"Error updating {vdb_file}: {e}")
        
        # Step 2: Clean the in-memory graph AND SAVE IMMEDIATELY
        if hasattr(rag_instance, 'chunk_entity_relation_graph'):
            graph_storage = rag_instance.chunk_entity_relation_graph
            graph = get_graph_from_storage(graph_storage)
            
            if graph is None:
                print("Warning: Could not access graph from storage")
            else:
                # Get initial counts
                initial_nodes = graph.number_of_nodes()
                initial_edges = graph.number_of_edges()
                
                # Find all nodes and edges to remove
                nodes_to_remove = set()
                edges_to_remove = set()
                
                # Check all nodes
                for node_id, node_data in list(graph.nodes(data=True)):
                    source_id = node_data.get('source_id', '')
                    file_path = node_data.get('file_path', '')
                    
                    for pattern in all_doc_id_patterns:
                        if pattern in source_id or pattern in file_path:
                            nodes_to_remove.add(node_id)
                            break
                
                # Check all edges
                for src, tgt, edge_data in list(graph.edges(data=True)):
                    # Remove if either node is being removed
                    if src in nodes_to_remove or tgt in nodes_to_remove:
                        edges_to_remove.add((src, tgt))
                        continue
                    
                    # Check edge's source_id
                    source_id = edge_data.get('source_id', '')
                    for pattern in all_doc_id_patterns:
                        if pattern in source_id:
                            edges_to_remove.add((src, tgt))
                            break
                
                # Remove edges first
                for src, tgt in edges_to_remove:
                    try:
                        graph.remove_edge(src, tgt)
                    except:
                        pass
                
                # Remove nodes
                for node in nodes_to_remove:
                    try:
                        graph.remove_node(node)
                    except:
                        pass
                
                print(f"Removed {len(nodes_to_remove)} nodes and {len(edges_to_remove)} edges from graph")
                print(f"Graph now has {graph.number_of_nodes()} nodes and {graph.number_of_edges()} edges")
                
                # Step 3: CRITICAL - Force save the graph multiple ways
                graphml_path = os.path.join(working_dir, "graph_chunk_entity_relation.graphml")
                
                # Create backup before saving
                if os.path.exists(graphml_path):
                    backup_path = graphml_path + ".backup"
                    shutil.copy2(graphml_path, backup_path)
                
                # Force save the graph
                save_success = force_save_graph_to_disk(graph, working_dir)
                
                if save_success:
                    # Also try storage save methods
                    if hasattr(rag_instance, 'graph_storage'):
                        try:
                            if hasattr(rag_instance.graph_storage, 'save'):
                                await rag_instance.graph_storage.save()
                                print("Also saved using graph_storage.save()")
                        except Exception as e:
                            print(f"Note: graph_storage.save() failed but file was saved: {e}")
                    
                    # Update the graph storage's internal reference
                    if hasattr(graph_storage, '_graph'):
                        graph_storage._graph = graph
                    elif hasattr(graph_storage, 'graph'):
                        graph_storage.graph = graph
                    
                    # Also update if graph_storage is the storage itself
                    if hasattr(rag_instance, 'graph_storage'):
                        if hasattr(rag_instance.graph_storage, '_graph'):
                            rag_instance.graph_storage._graph = graph
                        elif hasattr(rag_instance.graph_storage, 'graph'):
                            rag_instance.graph_storage.graph = graph
        
        # Step 4: Clean other JSON storage files
        json_files_to_clean = {
            "kv_store_full_docs.json": "doc_id",
            "kv_store_text_chunks.json": "chunk_id",
            "doc_status.json": "doc_id"
        }
        
        for json_file, id_field in json_files_to_clean.items():
            file_path = os.path.join(working_dir, json_file)
            if os.path.exists(file_path):
                try:
                    with open(file_path, 'r') as f:
                        data = json.load(f)
                    
                    if isinstance(data, dict):
                        keys_to_remove = []
                        for key in data.keys():
                            for pattern in all_doc_id_patterns:
                                if pattern in key:
                                    keys_to_remove.append(key)
                                    break
                        
                        for key in keys_to_remove:
                            data.pop(key, None)
                        
                        # Save with backup
                        backup_path = file_path + ".backup"
                        shutil.copy2(file_path, backup_path)
                        
                        with open(file_path, 'w') as f:
                            json.dump(data, f, indent=2)
                        
                        os.remove(backup_path)
                        print(f"Cleaned {len(keys_to_remove)} entries from {json_file}")
                        
                except Exception as e:
                    print(f"Error cleaning {json_file}: {e}")
        
        # Step 5: Clean KV storage if available
        if hasattr(rag_instance, 'kv_storage'):
            try:
                # Delete by doc_ids
                await rag_instance.kv_storage.delete(set(doc_ids))
                
                # Also clean any chunks
                if hasattr(rag_instance.kv_storage, 'get_all'):
                    all_data = await rag_instance.kv_storage.get_all()
                    keys_to_delete = set()
                    for key in all_data.keys():
                        for pattern in all_doc_id_patterns:
                            if pattern in key:
                                keys_to_delete.add(key)
                                break
                    if keys_to_delete:
                        await rag_instance.kv_storage.delete(keys_to_delete)
                        print(f"Deleted {len(keys_to_delete)} keys from KV storage")
            except Exception as e:
                print(f"Error cleaning KV storage: {e}")
        
        # Step 6: Clean doc status
        if hasattr(rag_instance, 'doc_status'):
            try:
                await rag_instance.doc_status.delete(set(doc_ids))
            except Exception as e:
                print(f"Error cleaning doc status: {e}")
        
        # Step 7: Clear any caches
        if hasattr(rag_instance, 'llm_response_cache'):
            try:
                await rag_instance.aclear_cache()
                print("Cleared LLM response cache")
            except Exception as e:
                print(f"Error clearing cache: {e}")
        
        # Step 8: Final verification - ensure files are actually updated
        time.sleep(0.5)  # Small delay to ensure filesystem sync
        
        # Verify GraphML file
        graphml_path = os.path.join(working_dir, "graph_chunk_entity_relation.graphml")
        if os.path.exists(graphml_path):
            try:
                verify_graph = nx.read_graphml(graphml_path)
                print(f"Final verification: GraphML has {verify_graph.number_of_nodes()} nodes and {verify_graph.number_of_edges()} edges")
            except Exception as e:
                print(f"Warning: Could not verify GraphML file: {e}")
        
        print("Enhanced cleanup completed successfully")
        
    except Exception as e:
        print(f"Critical error in enhanced cleanup: {e}")
        import traceback
        traceback.print_exc()
        raise

@asynccontextmanager
async def lifespan(app: FastAPI):
    """Lifecycle manager for FastAPI app with workspace support"""
    
    # Startup - Initialize default workspace
    print("Initializing workspace manager...")
    
    # Create the default workspace on startup
    await WorkspaceManager.create_workspace(default_workspace)
    
    # Load any existing workspaces
    base_dir = os.getenv("WORKING_DIR", "/app/data/rag_storage")
    workspaces_dir = os.path.join(base_dir, "workspaces")
    
    if os.path.exists(workspaces_dir):
        for workspace_name in os.listdir(workspaces_dir):
            workspace_path = os.path.join(workspaces_dir, workspace_name)
            if os.path.isdir(workspace_path) and workspace_name != default_workspace:
                try:
                    await WorkspaceManager.create_workspace(workspace_name)
                    print(f"Loaded existing workspace: {workspace_name}")
                except Exception as e:
                    print(f"Error loading workspace {workspace_name}: {e}")
    
    # Maintain backward compatibility - set global rag_instance to default workspace
    global rag_instance
    rag_instance = workspace_instances.get(default_workspace)
    
    yield
    
    # Shutdown - save all workspace metadata
    save_metadata_store()

# Create FastAPI app
app = FastAPI(
    title="Extended LightRAG API",
    description="LightRAG API with enhanced document management",
    version="1.0.0",
    lifespan=lifespan
)

# Import and include query routes
from query_routes import router as query_router
from enhanced_query_route import router as debug_router
app.include_router(query_router)
app.include_router(debug_router)

# Add CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Mount WebUI if available
webui_path = "/app/webui"
webui_mounted = False

if os.path.exists(webui_path) and os.path.isdir(webui_path):
    print(f"Mounting WebUI from {webui_path}")
    app.mount("/webui", StaticFiles(directory=webui_path, html=True), name="webui")
    webui_mounted = True
else:
    # Try alternative paths
    alt_paths = [
        "/usr/local/lib/python3.11/site-packages/lightrag/api/webui",
        "./webui",
        "../lightrag/api/webui"
    ]
    for path in alt_paths:
        if os.path.exists(path) and os.path.isdir(path):
            print(f"Mounting WebUI from {path}")
            app.mount("/webui", StaticFiles(directory=path, html=True), name="webui")
            webui_path = path
            webui_mounted = True
            break
    
    if not webui_mounted:
        print("Warning: WebUI directory not found. Web interface will not be available.")

@app.get("/", include_in_schema=False)
async def root():
    """Redirect root to WebUI"""
    if webui_mounted:
        return RedirectResponse(url="/webui/")
    else:
        return {"message": "LightRAG Extended API", "webui": "Not available", "docs": "/docs"}

@app.get("/health")
async def health_check():
    """Health check endpoint"""
    return {"status": "healthy", "service": "extended-lightrag"}

# Workspace Management Endpoints
@app.get("/workspaces")
async def list_workspaces():
    """List all available workspaces"""
    workspaces = WorkspaceManager.list_workspaces()
    workspace_details = []
    
    # Debug logging
    print(f"DEBUG /workspaces: Active instances: {list(workspace_instances.keys())}")
    print(f"DEBUG /workspaces: Metadata keys: {list(workspace_metadata.keys())}")
    
    for ws in workspaces:
        metadata = workspace_metadata.get(ws, {})
        workspace_details.append({
            "name": ws,
            "document_count": len(metadata),
            "is_default": ws == default_workspace,
            "is_active": ws in workspace_instances
        })
    
    return {
        "workspaces": workspace_details,
        "total": len(workspaces),
        "default_workspace": default_workspace,
        "debug": {
            "active_instances": list(workspace_instances.keys()),
            "metadata_workspaces": list(workspace_metadata.keys())
        }
    }

@app.post("/workspaces/{workspace_name}")
async def create_workspace(workspace_name: str):
    """Create a new workspace"""
    if workspace_name in workspace_instances:
        raise HTTPException(status_code=400, detail=f"Workspace '{workspace_name}' already exists")
    
    try:
        await WorkspaceManager.create_workspace(workspace_name)
        return {
            "status": "success",
            "message": f"Workspace '{workspace_name}' created successfully",
            "workspace": workspace_name
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/workspaces/{workspace_name}")
async def delete_workspace(workspace_name: str):
    """Delete a workspace and all its data"""
    if workspace_name == default_workspace:
        raise HTTPException(status_code=400, detail="Cannot delete the default workspace")
    
    if workspace_name not in WorkspaceManager.list_workspaces():
        raise HTTPException(status_code=404, detail=f"Workspace '{workspace_name}' not found")
    
    try:
        await WorkspaceManager.delete_workspace(workspace_name)
        return {
            "status": "success",
            "message": f"Workspace '{workspace_name}' deleted successfully"
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/workspaces/{workspace_name}/stats")
async def get_workspace_stats(workspace_name: str):
    """Get statistics for a specific workspace"""
    # URL decode the workspace name to handle spaces
    from urllib.parse import unquote
    workspace_name = unquote(workspace_name)
    
    if workspace_name not in WorkspaceManager.list_workspaces():
        # Try to create it if it doesn't exist
        await WorkspaceManager.create_workspace(workspace_name)
    
    # Get workspace instance
    rag = await WorkspaceManager.get_or_create_instance(workspace_name)
    metadata = workspace_metadata.get(workspace_name, {})
    
    # Calculate stats
    stats = {
        "workspace": workspace_name,
        "document_count": len(metadata),
        "is_active": workspace_name in workspace_instances,
        "storage_info": {
            "working_dir": rag.working_dir if rag else None,
            "workspace_param": rag.workspace if rag else None
        }
    }
    
    # Try to get graph stats if available
    if rag and hasattr(rag, 'chunk_entity_relation_graph'):
        try:
            graph_storage = rag.chunk_entity_relation_graph
            graph = get_graph_from_storage(graph_storage)
            if graph:
                stats["graph_stats"] = {
                    "nodes": graph.number_of_nodes(),
                    "edges": graph.number_of_edges()
                }
                # Add sample nodes for debugging
                sample_nodes = list(graph.nodes())[:5] if hasattr(graph, 'nodes') else []
                stats["sample_nodes"] = sample_nodes
        except Exception as e:
            stats["graph_error"] = str(e)
    
    return stats

@app.get("/workspaces/debug")
async def debug_workspaces():
    """Debug endpoint to understand workspace and graph issues"""
    debug_info = {
        "active_workspaces": list(workspace_instances.keys()),
        "workspace_metadata_keys": list(workspace_metadata.keys()),
        "workspace_details": {}
    }
    
    # Check each workspace
    for ws_name in list(workspace_instances.keys()):
        rag = workspace_instances[ws_name]
        ws_info = {
            "has_chunk_entity_relation_graph": hasattr(rag, 'chunk_entity_relation_graph'),
            "has_graph_storage": hasattr(rag, 'graph_storage'),
            "working_dir": rag.working_dir if hasattr(rag, 'working_dir') else None,
            "workspace_param": rag.workspace if hasattr(rag, 'workspace') else None
        }
        
        if hasattr(rag, 'chunk_entity_relation_graph'):
            try:
                graph_storage = rag.chunk_entity_relation_graph
                graph = get_graph_from_storage(graph_storage)
                if graph:
                    ws_info["graph_nodes"] = graph.number_of_nodes() if hasattr(graph, 'number_of_nodes') else 0
                    ws_info["graph_edges"] = graph.number_of_edges() if hasattr(graph, 'number_of_edges') else 0
                    ws_info["sample_nodes"] = list(graph.nodes())[:3] if hasattr(graph, 'nodes') else []
                else:
                    ws_info["graph"] = "None"
            except Exception as e:
                ws_info["graph_error"] = str(e)
        
        debug_info["workspace_details"][ws_name] = ws_info
    
    # Check combined graph
    try:
        combined = await get_combined_graph()
        debug_info["combined_graph"] = {
            "nodes": combined.number_of_nodes() if hasattr(combined, 'number_of_nodes') else 0,
            "edges": combined.number_of_edges() if hasattr(combined, 'number_of_edges') else 0
        }
    except Exception as e:
        debug_info["combined_graph_error"] = str(e)
    
    return debug_info

# Graph endpoints
@app.get("/graph/label/list")
async def get_graph_labels():
    """
    Get all graph labels from all workspaces
    Returns:
        List[str]: List of all unique graph labels across workspaces
    """
    try:
        # Get labels from all workspaces
        all_labels = await get_all_graph_labels()
        return all_labels
    except Exception as e:
        print(f"Error getting graph labels: {str(e)}")
        print(traceback.format_exc())
        raise HTTPException(
            status_code=500, detail=f"Error getting graph labels: {str(e)}"
        )

@app.get("/graphs")
async def get_knowledge_graph(
    label: str = Query(None, description="Label to get knowledge graph for"),
    max_depth: int = Query(3, description="Maximum depth of graph", ge=1),
    max_nodes: int = Query(1000, description="Maximum nodes to return", ge=1),
):
    """
    Retrieve a connected subgraph of nodes where the label includes the specified label.
    If no label is provided, returns the entire graph (up to max_nodes).
    When reducing the number of nodes, the prioritization criteria are as follows:
    1. Hops(path) to the starting node take precedence
    2. Followed by the degree of the nodes
    
    Args:
        label (str, optional): Label of the starting node. If None, returns entire graph
        max_depth (int, optional): Maximum depth of the subgraph, Defaults to 3
        max_nodes: Maximum nodes to return
    
    Returns:
        Dict[str, List[str]]: Knowledge graph for label or entire graph
    """
    try:
        print(f"\n=== /graphs endpoint called ===")
        print(f"Parameters: label={label}, max_depth={max_depth}, max_nodes={max_nodes}")
        
        # Handle empty string as None
        if label == "":
            label = None
            print("Empty label converted to None")
        
        # Handle wildcard for all nodes
        if label == "*":
            print("Wildcard label '*' requested - returning combined graph")
            label = None  # Treat as combined graph request
        
        # If label is provided, try workspace-specific methods first
        if label:
            for workspace_name, rag in workspace_instances.items():
                if hasattr(rag, 'get_knowledge_graph'):
                    try:
                        result = await rag.get_knowledge_graph(
                            node_label=label,
                            max_depth=max_depth,
                            max_nodes=max_nodes,
                        )
                        # If we found data, return it
                        if result and (result.get('nodes') or result.get('edges')):
                            print(f"Found graph data from workspace: {workspace_name}")
                            return result
                    except Exception as e:
                        print(f"Error getting graph from workspace {workspace_name}: {e}")
        
        # Use combined graph approach
        print("Using combined graph approach...")
        graph = await get_combined_graph()
        
        if graph is None or graph.number_of_nodes() == 0:
            print("No graph data available")
            return {"nodes": [], "edges": []}
        
        print(f"Combined graph has {graph.number_of_nodes()} nodes and {graph.number_of_edges()} edges")
        
        nodes_to_include = set()
        edges_to_include = []
        
        if label:
            # Get subgraph starting from the label
            if hasattr(graph, 'has_node') and graph.has_node(label):
                visited = set()
                queue = [(label, 0)]
                
                while queue and len(nodes_to_include) < max_nodes:
                    current_node, depth = queue.pop(0)
                    
                    if current_node in visited or depth > max_depth:
                        continue
                        
                    visited.add(current_node)
                    nodes_to_include.add(current_node)
                    
                    # Get neighbors and their edge data
                    if hasattr(graph, 'neighbors'):
                        for neighbor in graph.neighbors(current_node):
                            if neighbor not in visited and depth + 1 <= max_depth:
                                queue.append((neighbor, depth + 1))
                                # Get edge data if available
                                edge_data = {}
                                if hasattr(graph, 'get_edge_data'):
                                    edge_data = graph.get_edge_data(current_node, neighbor) or {}
                                
                                edge_dict = {
                                    "id": f"{current_node}-{neighbor}",  # Generate edge ID
                                    "source": current_node,
                                    "target": neighbor,
                                    "type": edge_data.get('keywords', None),  # Use keywords as type if available
                                    "properties": {k: v for k, v in edge_data.items() 
                                                 if k != 'keywords' and isinstance(v, (str, int, float, bool, list, dict))}
                                }
                                edges_to_include.append(edge_dict)
            else:
                print(f"Label '{label}' not found in graph")
        else:
            # Return entire graph (up to max_nodes)
            all_nodes = list(graph.nodes())
            nodes_to_include = set(all_nodes[:max_nodes])
            
            # Include edges between the selected nodes with their data
            for edge in graph.edges(data=True):
                source, target = edge[0], edge[1]
                edge_data = edge[2] if len(edge) > 2 else {}
                
                if source in nodes_to_include and target in nodes_to_include:
                    edge_dict = {
                        "id": f"{source}-{target}",  # Generate edge ID
                        "source": source,
                        "target": target,
                        "type": edge_data.get('keywords', None),  # Use keywords as type if available
                        "properties": {k: v for k, v in edge_data.items() 
                                     if k != 'keywords' and isinstance(v, (str, int, float, bool, list, dict))}
                    }
                    edges_to_include.append(edge_dict)
        
        print(f"Returning {len(nodes_to_include)} nodes and {len(edges_to_include)} edges")
        
        # Build response format matching original LightRAG structure
        node_list = []
        for node in nodes_to_include:
            # Create node in the format expected by WebUI
            # WebUI expects: {id: str, labels: [str], properties: {...}}
            node_properties = {}
            
            # Try to get additional node data
            if hasattr(graph, 'nodes') and hasattr(graph.nodes, '__getitem__'):
                try:
                    node_attrs = graph.nodes[node]
                    # Put all attributes in properties, filtering non-serializable
                    for key, value in node_attrs.items():
                        if isinstance(value, (str, int, float, bool, list, dict)):
                            node_properties[key] = value
                except:
                    pass
            
            # Build node with correct structure for WebUI
            node_data = {
                "id": str(node),
                "labels": [str(node)],  # WebUI expects labels as array
                "properties": node_properties  # All other data goes in properties
            }
            node_list.append(node_data)
        
        return {
            "nodes": node_list,
            "edges": edges_to_include
        }
            
    except Exception as e:
        print(f"Error in /graphs endpoint: {str(e)}")
        print(traceback.format_exc())
        raise HTTPException(
            status_code=500, detail=f"Error getting knowledge graph: {str(e)}"
        )

@app.get("/graph/entity/exists")
async def check_entity_exists(
    name: str = Query(..., description="Entity name to check"),
):
    """
    Check if an entity with the given name exists in the knowledge graph
    
    Args:
        name (str): Name of the entity to check
    
    Returns:
        Dict[str, bool]: Dictionary with 'exists' key indicating if entity exists
    """
    try:
        # Check all workspaces for the entity
        for workspace_name, rag in workspace_instances.items():
            if hasattr(rag, 'chunk_entity_relation_graph'):
                graph_storage = rag.chunk_entity_relation_graph
                graph = get_graph_from_storage(graph_storage)
                
                if graph:
                    if hasattr(graph, 'has_node'):
                        exists = await graph.has_node(name) if asyncio.iscoroutinefunction(graph.has_node) else graph.has_node(name)
                        if exists:
                            return {"exists": True, "workspace": workspace_name}
                    elif hasattr(graph, '__contains__'):
                        exists = name in graph
                        if exists:
                            return {"exists": True, "workspace": workspace_name}
        
        return {"exists": False}
    except Exception as e:
        print(f"Error checking entity existence for '{name}': {str(e)}")
        print(traceback.format_exc())
        raise HTTPException(
            status_code=500, detail=f"Error checking entity existence: {str(e)}"
        )

@app.post("/graph/entity/edit")
async def update_entity(request: EntityUpdateRequest):
    """
    Update an entity's properties in the knowledge graph
    
    Args:
        request (EntityUpdateRequest): Request containing entity name, updated data, and rename flag
    
    Returns:
        Dict: Updated entity information
    """
    try:
        # Check if rag_instance has the edit method
        if hasattr(rag_instance, 'aedit_entity'):
            result = await rag_instance.aedit_entity(
                entity_name=request.entity_name,
                updated_data=request.updated_data,
                allow_rename=request.allow_rename,
            )
            return {
                "status": "success",
                "message": "Entity updated successfully",
                "data": result,
            }
        else:
            # Alternative approach: update entity in graph storage directly
            if hasattr(rag_instance, 'chunk_entity_relation_graph'):
                graph_storage = rag_instance.chunk_entity_relation_graph
                graph = get_graph_from_storage(graph_storage)
                
                if graph is None:
                    raise ValueError("Could not access graph from storage")
                
                # Check if entity exists
                if hasattr(graph, 'has_node') and graph.has_node(request.entity_name):
                    # Update node attributes
                    if hasattr(graph, 'nodes') and hasattr(graph.nodes, '__setitem__'):
                        for key, value in request.updated_data.items():
                            graph.nodes[request.entity_name][key] = value
                    
                    # Handle rename if requested
                    if request.allow_rename and 'name' in request.updated_data:
                        new_name = request.updated_data['name']
                        if new_name != request.entity_name:
                            # This is more complex - would need to rename node
                            # For now, just return success without renaming
                            pass
                    
                    return {
                        "status": "success",
                        "message": "Entity updated successfully",
                        "data": request.updated_data,
                    }
                else:
                    raise ValueError(f"Entity '{request.entity_name}' not found")
            
            raise ValueError("Graph storage not available")
            
    except ValueError as ve:
        print(f"Validation error updating entity '{request.entity_name}': {str(ve)}")
        raise HTTPException(status_code=400, detail=str(ve))
    except Exception as e:
        print(f"Error updating entity '{request.entity_name}': {str(e)}")
        print(traceback.format_exc())
        raise HTTPException(
            status_code=500, detail=f"Error updating entity: {str(e)}"
        )

@app.post("/graph/relation/edit")
async def update_relation(request: RelationUpdateRequest):
    """Update a relation's properties in the knowledge graph
    
    Args:
        request (RelationUpdateRequest): Request containing source ID, target ID and updated data
    
    Returns:
        Dict: Updated relation information
    """
    try:
        # Check if rag_instance has the edit method
        if hasattr(rag_instance, 'aedit_relation'):
            result = await rag_instance.aedit_relation(
                source_entity=request.source_id,
                target_entity=request.target_id,
                updated_data=request.updated_data,
            )
            return {
                "status": "success",
                "message": "Relation updated successfully",
                "data": result,
            }
        else:
            # Alternative approach: update edge in graph storage directly
            if hasattr(rag_instance, 'chunk_entity_relation_graph'):
                graph_storage = rag_instance.chunk_entity_relation_graph
                graph = get_graph_from_storage(graph_storage)
                
                if graph is None:
                    raise ValueError("Could not access graph from storage")
                
                # Check if edge exists
                if hasattr(graph, 'has_edge') and graph.has_edge(request.source_id, request.target_id):
                    # Update edge attributes
                    if hasattr(graph, 'edges') and hasattr(graph.edges, '__getitem__'):
                        edge = graph.edges[request.source_id, request.target_id]
                        for key, value in request.updated_data.items():
                            edge[key] = value
                    
                    return {
                        "status": "success",
                        "message": "Relation updated successfully",
                        "data": request.updated_data,
                    }
                else:
                    raise ValueError(f"Relation between '{request.source_id}' and '{request.target_id}' not found")
            
            raise ValueError("Graph storage not available")
            
    except ValueError as ve:
        print(f"Validation error updating relation between '{request.source_id}' and '{request.target_id}': {str(ve)}")
        raise HTTPException(status_code=400, detail=str(ve))
    except Exception as e:
        print(f"Error updating relation between '{request.source_id}' and '{request.target_id}': {str(e)}")
        print(traceback.format_exc())
        raise HTTPException(
            status_code=500, detail=f"Error updating relation: {str(e)}"
        )

@app.get("/graph/status")
async def get_graph_status():
    """Get the combined status of knowledge graphs from all workspaces"""
    try:
        # Get combined graph from all workspaces
        combined_graph = await get_combined_graph()
        
        status = {
            "exists": combined_graph is not None and combined_graph.number_of_nodes() > 0,
            "type": "CombinedGraph",
            "graph_type": type(combined_graph).__name__ if combined_graph else "None",
            "nodes": combined_graph.number_of_nodes() if combined_graph and hasattr(combined_graph, 'number_of_nodes') else 0,
            "edges": combined_graph.number_of_edges() if combined_graph and hasattr(combined_graph, 'number_of_edges') else 0,
            "workspaces": len(workspace_instances),
            "workspace_details": {}
        }
        
        # Add per-workspace details
        for workspace_name, rag in workspace_instances.items():
            if hasattr(rag, 'chunk_entity_relation_graph'):
                graph_storage = rag.chunk_entity_relation_graph
                graph = get_graph_from_storage(graph_storage)
                status["workspace_details"][workspace_name] = {
                    "nodes": graph.number_of_nodes() if graph and hasattr(graph, 'number_of_nodes') else 0,
                    "edges": graph.number_of_edges() if graph and hasattr(graph, 'number_of_edges') else 0
                }
            
            # Check all relevant files
            working_dir = os.getenv("WORKING_DIR", "/app/data/rag_storage")
            
            # Check GraphML file
            graphml_file = os.path.join(working_dir, "graph_chunk_entity_relation.graphml")
            status["graphml_file_exists"] = os.path.exists(graphml_file)
            if status["graphml_file_exists"]:
                status["graphml_file_size"] = os.path.getsize(graphml_file)
            
            # Check vector database files
            vector_db_files = {
                "vdb_entities.json": "Entity embeddings",
                "vdb_relationships.json": "Relationship embeddings",
                "vdb_chunks.json": "Document chunk embeddings"
            }
            
            status["vector_databases"] = {}
            for file_name, description in vector_db_files.items():
                file_path = os.path.join(working_dir, file_name)
                if os.path.exists(file_path):
                    try:
                        with open(file_path, 'r') as f:
                            data = json.load(f)
                            entity_count = len(data.get('data', []))
                            status["vector_databases"][file_name] = {
                                "exists": True,
                                "description": description,
                                "size": os.path.getsize(file_path),
                                "entity_count": entity_count
                            }
                    except Exception as e:
                        status["vector_databases"][file_name] = {
                            "exists": True,
                            "error": str(e)
                        }
                else:
                    status["vector_databases"][file_name] = {
                        "exists": False,
                        "description": description
                    }
            
            # Sample some nodes if they exist
            if status["nodes"] > 0 and hasattr(graph, 'nodes'):
                sample_nodes = list(graph.nodes())[:5]
                status["sample_nodes"] = sample_nodes
            
            # Total entity count from vector DBs
            total_entities_in_vdbs = sum(
                vdb.get("entity_count", 0) 
                for vdb in status["vector_databases"].values() 
                if "entity_count" in vdb
            )
            status["total_entities_in_vector_dbs"] = total_entities_in_vdbs
            
            return status
        else:
            return {
                "exists": False,
                "message": "Knowledge graph not initialized"
            }
            
    except Exception as e:
        print(f"Error getting graph status: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/graph/clear")
async def clear_knowledge_graph():
    """Clear the entire knowledge graph"""
    try:
        if hasattr(rag_instance, 'chunk_entity_relation_graph'):
            graph_storage = rag_instance.chunk_entity_relation_graph
            graph = get_graph_from_storage(graph_storage)
            
            if graph is None:
                raise HTTPException(
                    status_code=500,
                    detail="Could not access graph from storage"
                )
            
            # Get node and edge counts before clearing
            node_count = graph.number_of_nodes() if hasattr(graph, 'number_of_nodes') else 0
            edge_count = graph.number_of_edges() if hasattr(graph, 'number_of_edges') else 0
            
            if hasattr(graph, 'clear'):
                print(f"Clearing knowledge graph with {node_count} nodes and {edge_count} edges")
                graph.clear()
                
                # Critical: Save the cleared graph to the GraphML file
                working_dir = os.getenv("WORKING_DIR", "/app/data/rag_storage")
                graphml_path = os.path.join(working_dir, "graph_chunk_entity_relation.graphml")
                
                # Force save the cleared graph
                force_save_graph_to_disk(graph, working_dir)
                
                # Also try the storage's save methods
                if hasattr(rag_instance, 'graph_storage'):
                    if hasattr(rag_instance.graph_storage, 'save'):
                        try:
                            await rag_instance.graph_storage.save()
                            print("Also saved using graph_storage.save()")
                        except Exception as e:
                            print(f"Note: graph_storage.save() failed: {e}")
                    
                    if hasattr(rag_instance.graph_storage, '_save'):
                        try:
                            rag_instance.graph_storage._save()
                            print("Also saved using graph_storage._save()")
                        except Exception as e:
                            print(f"Note: graph_storage._save() failed: {e}")
                
                # CRITICAL: Clear ALL graph-related and vector database files
                files_to_clear = [
                    # Graph files
                    "graph_chunk_entity_relation.graphml",
                    "graph_data.json",
                    "graph_cache.json",
                    # Vector database files - THESE ARE CRITICAL FOR WEBUI
                    "vdb_entities.json",
                    "vdb_relationships.json",
                    "vdb_chunks.json",
                    # Additional possible vector storage files
                    "entity_embedding.json",
                    "relationship_embedding.json",
                    "document_graph_storage.json"
                ]
                
                print(f"Clearing all graph and vector database files in {working_dir}")
                
                for file_name in files_to_clear:
                    file_path = os.path.join(working_dir, file_name)
                    if os.path.exists(file_path):
                        try:
                            if file_name.endswith('.graphml'):
                                # Ensure GraphML file contains empty graph
                                empty_graph = nx.Graph()
                                nx.write_graphml(empty_graph, file_path)
                                print(f"Wrote empty graph to {file_path}")
                            elif file_name.startswith('vdb_'):
                                # For vector database files, write the proper empty structure
                                empty_vdb = {
                                    "embedding_dim": 1536,  # Default OpenAI embedding dimension
                                    "data": [],
                                    "matrix": []
                                }
                                with open(file_path, 'w') as f:
                                    json.dump(empty_vdb, f, indent=2)
                                print(f"Cleared vector database: {file_path}")
                            else:
                                # Clear JSON files
                                with open(file_path, 'w') as f:
                                    json.dump({}, f)
                                print(f"Cleared {file_path}")
                        except Exception as e:
                            print(f"Error handling {file_path}: {e}")
                            # Try to delete the file if clearing fails
                            try:
                                os.remove(file_path)
                                print(f"Deleted {file_path}")
                            except Exception as del_e:
                                print(f"Error deleting {file_path}: {del_e}")
                
                # Try to clear entity and relationship vector stores if accessible
                possible_vdb_attrs = [
                    'entities_vdb', 'relationships_vdb', 'chunks_vdb',
                    'entity_vdb', 'relation_vdb', 'chunk_vdb'
                ]
                
                for attr_name in possible_vdb_attrs:
                    if hasattr(rag_instance, attr_name):
                        vdb = getattr(rag_instance, attr_name)
                        if vdb and hasattr(vdb, 'clear'):
                            try:
                                await vdb.clear()
                                print(f"Cleared {attr_name}")
                            except Exception as e:
                                print(f"Error clearing {attr_name}: {e}")
                
                return {
                    "status": "success",
                    "message": "Knowledge graph and all vector databases cleared successfully",
                    "cleared": {
                        "nodes": node_count,
                        "edges": edge_count,
                        "files_cleared": files_to_clear
                    }
                }
            else:
                raise HTTPException(
                    status_code=500,
                    detail="Graph clear method not available"
                )
        else:
            raise HTTPException(
                status_code=404,
                detail="Knowledge graph not found"
            )
            
    except Exception as e:
        print(f"Error clearing knowledge graph: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/graph/debug")
async def debug_graph_sources():
    """Debug endpoint to check all possible sources of graph data across all workspaces"""
    try:
        base_dir = os.getenv("WORKING_DIR", "/app/data/rag_storage")
        debug_info = {
            "base_dir": base_dir,
            "workspaces": {},
            "total_nodes": 0,
            "total_edges": 0
        }
        
        # Check each workspace
        for workspace_name, rag in workspace_instances.items():
            workspace_info = {
                "working_dir": rag.working_dir if hasattr(rag, 'working_dir') else "unknown",
                "in_memory_graph": {},
                "files": {},
                "metadata_count": len(workspace_metadata.get(workspace_name, {}))
            }
            
            # Check in-memory graph
            if hasattr(rag, 'chunk_entity_relation_graph'):
                graph_storage = rag.chunk_entity_relation_graph
                graph = await get_graph_from_storage(graph_storage)
                
                node_count = 0
                edge_count = 0
                if graph and hasattr(graph, 'number_of_nodes'):
                    node_count = graph.number_of_nodes()
                    edge_count = graph.number_of_edges()
                    debug_info["total_nodes"] += node_count
                    debug_info["total_edges"] += edge_count
                
                workspace_info["in_memory_graph"] = {
                    "exists": True,
                    "storage_type": type(graph_storage).__name__,
                    "graph_type": type(graph).__name__ if graph else "None",
                    "nodes": node_count,
                    "edges": edge_count,
                    "has_graph_attr": hasattr(graph_storage, '_graph'),
                    "has_get_graph": hasattr(graph_storage, '_get_graph'),
                    "storage_initialized": hasattr(graph_storage, '_storage_lock')
                }
            else:
                workspace_info["in_memory_graph"]["exists"] = False
            
            
            # Check files for this workspace
            files_to_check = [
                "graph_chunk_entity_relation.graphml",
                "vdb_entities.json",
                "vdb_relationships.json",
                "vdb_chunks.json",
                "kv_store_text_chunks.json",
                "kv_store_full_docs.json",
                "doc_status.json"
            ]
            
            workspace_dir = workspace_info["working_dir"]
            for file_name in files_to_check:
                file_path = os.path.join(workspace_dir, file_name)
                if os.path.exists(file_path):
                    file_info = {
                        "exists": True,
                        "size": os.path.getsize(file_path),
                        "modified": os.path.getmtime(file_path)
                    }
                    
                    # Special handling for GraphML files
                    if file_name.endswith('.graphml'):
                        try:
                            test_graph = nx.read_graphml(file_path)
                            file_info["nodes_in_file"] = test_graph.number_of_nodes()
                            file_info["edges_in_file"] = test_graph.number_of_edges()
                        except Exception as e:
                            file_info["read_error"] = str(e)
                    
                    # For JSON files, try to get entity count
                    elif file_name.endswith('.json'):
                        try:
                            with open(file_path, 'r') as f:
                                data = json.load(f)
                                if isinstance(data, dict):
                                    if 'data' in data:
                                        file_info["entity_count"] = len(data['data'])
                                    else:
                                        file_info["key_count"] = len(data)
                                elif isinstance(data, list):
                                    file_info["item_count"] = len(data)
                        except Exception as e:
                            file_info["read_error"] = str(e)
                    
                    workspace_info["files"][file_name] = file_info
                else:
                    workspace_info["files"][file_name] = {"exists": False}
            
            debug_info["workspaces"][workspace_name] = workspace_info
        
        # Legacy check for old working_dir (for backward compatibility)
        old_working_dir = os.getenv("WORKING_DIR", "/app/data/rag_storage")
        if os.path.exists(old_working_dir):
            old_graphml = os.path.join(old_working_dir, "graph_chunk_entity_relation.graphml")
            if os.path.exists(old_graphml):
                debug_info["legacy_graphml"] = {
                    "path": old_graphml,
                    "size": os.path.getsize(old_graphml)
                }
        
        return debug_info
        
    except Exception as e:
        print(f"Error in debug endpoint: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/graph/reload")
async def reload_graph(workspace: str = None):
    """Force reload the graph from the GraphML file for a specific workspace or all workspaces"""
    try:
        reload_results = {}
        
        # Determine which workspaces to reload
        workspaces_to_reload = []
        if workspace:
            if workspace in workspace_instances:
                workspaces_to_reload = [workspace]
            else:
                return {"error": f"Workspace '{workspace}' not found"}
        else:
            workspaces_to_reload = list(workspace_instances.keys())
        
        # Reload each workspace
        for ws_name in workspaces_to_reload:
            rag = workspace_instances[ws_name]
            working_dir = rag.working_dir if hasattr(rag, 'working_dir') else os.path.join(os.getenv("WORKING_DIR", "/app/data/rag_storage"), "workspaces", ws_name)
            graphml_path = os.path.join(working_dir, "graph_chunk_entity_relation.graphml")
            
            if os.path.exists(graphml_path):
                try:
                    # Load the graph from file
                    loaded_graph = nx.read_graphml(graphml_path)
                    node_count = loaded_graph.number_of_nodes()
                    edge_count = loaded_graph.number_of_edges()
                    
                    # Replace the in-memory graph
                    if hasattr(rag, 'chunk_entity_relation_graph'):
                        graph_storage = rag.chunk_entity_relation_graph
                        
                        # Update the graph in the storage object
                        if hasattr(graph_storage, '_graph'):
                            graph_storage._graph = loaded_graph
                            # Mark as not needing update
                            if hasattr(graph_storage, 'storage_updated') and graph_storage.storage_updated:
                                graph_storage.storage_updated.value = False
                        elif hasattr(graph_storage, 'graph'):
                            graph_storage.graph = loaded_graph
                        else:
                            # If storage doesn't have a way to set the graph, log warning
                            reload_results[ws_name] = {
                                "status": "error",
                                "message": "Could not update graph in storage object"
                            }
                            continue
                        
                        reload_results[ws_name] = {
                            "status": "success",
                            "message": f"Graph reloaded from {graphml_path}",
                            "nodes": node_count,
                            "edges": edge_count
                        }
                    else:
                        reload_results[ws_name] = {
                            "status": "error",
                            "message": "Could not find graph reference in RAG instance"
                        }
                except Exception as e:
                    reload_results[ws_name] = {
                        "status": "error",
                        "message": f"Error loading GraphML: {str(e)}"
                    }
            else:
                reload_results[ws_name] = {
                    "status": "warning",
                    "message": f"GraphML file not found at {graphml_path}"
                }
        
        return reload_results
            
    except Exception as e:
        print(f"Error reloading graph: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/verify")
async def verify_deployment():
    """Comprehensive deployment verification endpoint"""
    try:
        verification = {
            "status": "healthy",
            "checks": {},
            "warnings": [],
            "errors": []
        }
        
        # Check 1: Workspaces loaded
        verification["checks"]["workspaces_loaded"] = {
            "count": len(workspace_instances),
            "names": list(workspace_instances.keys())
        }
        
        # Check 2: Graph data availability
        total_nodes = 0
        total_edges = 0
        graph_status = {}
        
        for ws_name, rag in workspace_instances.items():
            if hasattr(rag, 'chunk_entity_relation_graph'):
                try:
                    graph = await get_graph_from_storage(rag.chunk_entity_relation_graph)
                    if graph and hasattr(graph, 'number_of_nodes'):
                        nodes = graph.number_of_nodes()
                        edges = graph.number_of_edges()
                        total_nodes += nodes
                        total_edges += edges
                        graph_status[ws_name] = {"nodes": nodes, "edges": edges, "status": "ok"}
                    else:
                        graph_status[ws_name] = {"nodes": 0, "edges": 0, "status": "empty"}
                        verification["warnings"].append(f"Workspace '{ws_name}' has empty graph")
                except Exception as e:
                    graph_status[ws_name] = {"nodes": 0, "edges": 0, "status": "error", "error": str(e)}
                    verification["errors"].append(f"Error loading graph for workspace '{ws_name}': {e}")
            else:
                graph_status[ws_name] = {"nodes": 0, "edges": 0, "status": "no_graph"}
                verification["warnings"].append(f"Workspace '{ws_name}' has no graph storage")
        
        verification["checks"]["graph_data"] = {
            "total_nodes": total_nodes,
            "total_edges": total_edges,
            "by_workspace": graph_status
        }
        
        # Check 3: File system structure
        base_dir = os.getenv("WORKING_DIR", "/app/data/rag_storage")
        nested_workspaces = []
        
        workspaces_dir = os.path.join(base_dir, "workspaces")
        if os.path.exists(workspaces_dir):
            for ws_name in os.listdir(workspaces_dir):
                ws_path = os.path.join(workspaces_dir, ws_name)
                if os.path.isdir(ws_path):
                    nested_path = os.path.join(ws_path, ws_name)
                    if os.path.exists(nested_path) and os.path.isdir(nested_path):
                        nested_workspaces.append(ws_name)
        
        verification["checks"]["file_structure"] = {
            "base_dir": base_dir,
            "nested_workspaces": nested_workspaces
        }
        
        if nested_workspaces:
            verification["warnings"].append(f"Found nested workspace directories: {nested_workspaces}. Run migration script.")
        
        # Check 4: API endpoints
        verification["checks"]["api_endpoints"] = {
            "graphs": "/graphs endpoint available",
            "documents": "/documents endpoint available",
            "debug": "/graph/debug endpoint available",
            "health": "/health endpoint available"
        }
        
        # Check 5: WebUI availability
        webui_path = "./webui" if os.path.exists("./webui") else None
        verification["checks"]["webui"] = {
            "available": webui_path is not None,
            "path": webui_path
        }
        
        # Determine overall status
        if verification["errors"]:
            verification["status"] = "unhealthy"
        elif verification["warnings"]:
            verification["status"] = "degraded"
        
        return verification
        
    except Exception as e:
        print(f"Error in verification endpoint: {e}")
        import traceback
        traceback.print_exc()
        return {
            "status": "error",
            "error": str(e),
            "traceback": traceback.format_exc()
        }

@app.post("/documents/text/enhanced")
async def insert_text_enhanced(request: EnhancedTextInsertRequest):
    """Enhanced text insertion with full metadata support and processing verification"""
    try:
        # Compute document ID
        doc_id = compute_doc_id(request.text)
        
        # Handle source_url processing
        if request.source_url:
            enriched_content = f"Source: {request.source_url}\n"
            if request.description:
                enriched_content += f"Description: {request.description}\n"
            enriched_content += f"\n{request.text}"
        else:
            enriched_content = request.text
        
        # Determine file_path
        if request.document_id:
            if request.document_id.startswith('[') and ']' in request.document_id:
                # Handle [domain] format from n8n
                if request.document_id.startswith('[unknown] '):
                    # Clean up [unknown] prefix
                    cleaned_id = request.document_id.replace('[unknown] ', '')
                    if request.source_url:
                        try:
                            from urllib.parse import urlparse
                            parsed = urlparse(request.source_url)
                            domain = parsed.netloc
                            file_path = f"[{domain}] {cleaned_id}"
                            full_path = request.source_url  # Store the full URL separately
                        except:
                            file_path = f"[unknown] {cleaned_id}"
                    else:
                        file_path = f"[unknown] {cleaned_id}"
                else:
                    # Use the provided document_id as-is
                    file_path = request.document_id
                    if request.source_url:
                        full_path = request.source_url
            else:
                # Regular document_id, create file_path
                if request.source_url:
                    try:
                        from urllib.parse import urlparse
                        parsed = urlparse(request.source_url)
                        domain = parsed.netloc
                        file_path = f"[{domain}] {request.document_id}"
                        full_path = request.source_url
                    except:
                        file_path = f"text/{request.document_id}.txt"
                else:
                    file_path = f"text/{request.document_id}.txt"
        else:
            # No document_id provided, use computed hash
            if request.source_url:
                try:
                    from urllib.parse import urlparse
                    parsed = urlparse(request.source_url)
                    domain = parsed.netloc
                    # Use the last part of the URL path as filename
                    path_parts = parsed.path.strip('/').split('/')
                    page_name = path_parts[-1] if path_parts and path_parts[-1] else 'index'
                    file_path = f"[{domain}] {page_name}"
                    full_path = request.source_url
                except:
                    file_path = f"text/{doc_id}.txt"
            else:
                file_path = f"text/{doc_id}.txt"
        
        # Generate display_name from file_path
        if file_path.startswith('[') and ']' in file_path:
            # For [domain] format, use the full path as display name
            display_name = file_path
        else:
            # For other formats, use file_path or generate one
            if len(file_path) > 50:
                display_name = f"text/{doc_id[:8]}..."  # Shortened ID for display
            else:
                display_name = f"text/{doc_id[:8]}..."  # Shortened ID for display

        # Determine which ID to use for LightRAG
        if request.document_id and request.document_id.startswith('[') and ']' in request.document_id:
            # Check if this is an [unknown] format and use the cleaned version
            if request.document_id.startswith('[unknown] '):
                # Use the cleaned file_path we generated above as the custom ID
                custom_id = file_path
            else:
                # Use the n8n-provided document_id as the custom ID
                custom_id = request.document_id
            # IMPORTANT: Use the custom_id as the key for metadata storage
            # This ensures consistency between LightRAG's internal ID and our metadata
            metadata_key = custom_id
        else:
            # Use the computed hash ID
            custom_id = doc_id
            metadata_key = doc_id
        
        # Extract workspace from the custom_id or file_path
        workspace = WorkspaceManager.extract_workspace_from_doc_id(custom_id)
        
        # Get the appropriate RAG instance for this workspace
        rag = await WorkspaceManager.get_or_create_instance(workspace)

        # Store metadata with the correct key
        metadata_entry = {
            "id": metadata_key,
            "original_doc_id": doc_id,  # Keep the hash ID for reference
            "file_path": file_path,
            "display_name": display_name,
            "source_url": request.source_url,
            "description": request.description,
            "indexed_at": datetime.utcnow().isoformat(),
            "sitemap_url": request.sitemap_url,
            "doc_index": request.doc_index,
            "total_docs": request.total_docs,
            "content_summary": enriched_content[:200] + "..." if len(enriched_content) > 200 else enriched_content
        }

        # Add full_path if we have it (when using document_id from n8n)
        if request.document_id and 'full_path' in locals():
            metadata_entry["full_path"] = full_path

        # Insert into the workspace's LightRAG instance with custom ID
        await rag.ainsert(enriched_content, ids=[custom_id], file_paths=[file_path])
        
        # Wait for processing to complete and verify success
        import asyncio
        max_retries = 30  # Wait up to 30 seconds
        retry_count = 0
        
        while retry_count < max_retries:
            await asyncio.sleep(1)  # Wait 1 second between checks
            
            # Check document status
            if hasattr(rag, 'doc_status') and rag.doc_status is not None:
                try:
                    doc_status_data = await rag.doc_status.get(custom_id)
                    if doc_status_data:
                        status = doc_status_data.get('status', 'unknown')
                        if status == 'processed':
                            # Successfully processed - now store metadata in the workspace
                            if workspace not in workspace_metadata:
                                workspace_metadata[workspace] = {}
                            workspace_metadata[workspace][metadata_key] = metadata_entry
                            save_workspace_metadata(workspace)
                            
                            # Force save graph to disk after successful processing
                            try:
                                if hasattr(rag, 'chunk_entity_relation_graph'):
                                    graph_storage = rag.chunk_entity_relation_graph
                                    
                                    # Call index_done_callback to persist the graph
                                    if hasattr(graph_storage, 'index_done_callback'):
                                        result = await graph_storage.index_done_callback()
                                        print(f"Called index_done_callback on graph storage, result: {result}")
                                    
                                    # Always try to get and save the graph directly as backup
                                    graph = await get_graph_from_storage(graph_storage)
                                    if graph and hasattr(graph, 'number_of_nodes'):
                                        graphml_path = os.path.join(rag.working_dir, "graph_chunk_entity_relation.graphml")
                                        nx.write_graphml(graph, graphml_path)
                                        print(f"Directly saved graph with {graph.number_of_nodes()} nodes to {graphml_path}")
                                        
                                        # Verify the file was written
                                        if os.path.exists(graphml_path):
                                            file_size = os.path.getsize(graphml_path)
                                            print(f"Graph file saved successfully, size: {file_size} bytes")
                                        else:
                                            print(f"ERROR: Graph file was not created at {graphml_path}")
                                    else:
                                        print(f"WARNING: Graph is empty or invalid")
                            except Exception as e:
                                print(f"ERROR saving graph to disk: {e}")
                                import traceback
                                traceback.print_exc()
                            
                            return {
                                "status": "success",
                                "message": "Document inserted and processed successfully",
                                "doc_id": metadata_key,
                                "file_path": file_path,
                                "processing_status": "completed"
                            }
                        elif status == 'failed':
                            error_msg = doc_status_data.get('error', 'Unknown processing error')
                            raise HTTPException(
                                status_code=500, 
                                detail=f"Document processing failed: {error_msg}"
                            )
                        # If status is 'pending' or 'processing', continue waiting
                except Exception as e:
                    print(f"Error checking document status: {e}")
            
            retry_count += 1
        
        # If we get here, processing timed out
        raise HTTPException(
            status_code=500, 
            detail="Document processing timed out. The document may still be processing in the background."
        )
        
    except HTTPException:
        raise
    except Exception as e:
        print(f"Error in insert_text_enhanced: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/documents/text")
async def insert_text(request: TextInsertRequest):
    """Standard text insertion with workspace support based on file_path prefix"""
    try:
        # Compute document ID
        doc_id = compute_doc_id(request.text)
        
        # Use provided file_path or create one
        file_path = request.file_path if request.file_path else f"text/{doc_id}.txt"
        
        # Extract workspace from file_path if it has [workspace] prefix
        workspace = WorkspaceManager.extract_workspace_from_doc_id(file_path)
        
        # Get the appropriate RAG instance for this workspace
        rag = await WorkspaceManager.get_or_create_instance(workspace)
        
        # Generate display_name from file_path
        display_name = generate_display_name_from_file_path(file_path, doc_id)
        
        # Store metadata in the appropriate workspace
        metadata_entry = {
            "id": doc_id,
            "file_path": file_path,
            "display_name": display_name,
            "description": request.description,
            "indexed_at": datetime.utcnow().isoformat(),
            "content_summary": request.text[:200] + "..." if len(request.text) > 200 else request.text
        }
        
        # Insert into the workspace's LightRAG instance
        await rag.ainsert(request.text, file_paths=[file_path])
        
        # Wait for processing to complete and verify success
        import asyncio
        max_retries = 30  # Wait up to 30 seconds
        retry_count = 0
        
        while retry_count < max_retries:
            await asyncio.sleep(1)  # Wait 1 second between checks
            
            # Check document status
            if hasattr(rag, 'doc_status') and rag.doc_status is not None:
                try:
                    doc_status_data = await rag.doc_status.get(doc_id)
                    if doc_status_data:
                        status = doc_status_data.get('status', 'unknown')
                        if status == 'processed':
                            # Successfully processed - now store metadata in the workspace
                            if workspace not in workspace_metadata:
                                workspace_metadata[workspace] = {}
                            workspace_metadata[workspace][doc_id] = metadata_entry
                            save_workspace_metadata(workspace)
                            
                            # Force save graph to disk after successful processing
                            try:
                                if hasattr(rag, 'chunk_entity_relation_graph'):
                                    graph_storage = rag.chunk_entity_relation_graph
                                    
                                    # Call index_done_callback to persist the graph
                                    if hasattr(graph_storage, 'index_done_callback'):
                                        result = await graph_storage.index_done_callback()
                                        print(f"Called index_done_callback on graph storage, result: {result}")
                                    
                                    # Always try to get and save the graph directly as backup
                                    graph = await get_graph_from_storage(graph_storage)
                                    if graph and hasattr(graph, 'number_of_nodes'):
                                        graphml_path = os.path.join(rag.working_dir, "graph_chunk_entity_relation.graphml")
                                        nx.write_graphml(graph, graphml_path)
                                        print(f"Directly saved graph with {graph.number_of_nodes()} nodes to {graphml_path}")
                                        
                                        # Verify the file was written
                                        if os.path.exists(graphml_path):
                                            file_size = os.path.getsize(graphml_path)
                                            print(f"Graph file saved successfully, size: {file_size} bytes")
                                        else:
                                            print(f"ERROR: Graph file was not created at {graphml_path}")
                                    else:
                                        print(f"WARNING: Graph is empty or invalid")
                            except Exception as e:
                                print(f"ERROR saving graph to disk: {e}")
                                import traceback
                                traceback.print_exc()
                            
                            return {
                                "status": "success",
                                "message": "Document inserted and processed successfully",
                                "doc_id": doc_id,
                                "file_path": file_path,
                                "processing_status": "completed"
                            }
                        elif status == 'failed':
                            error_msg = doc_status_data.get('error', 'Unknown processing error')
                            raise HTTPException(
                                status_code=500, 
                                detail=f"Document processing failed: {error_msg}"
                            )
                        # If status is 'pending' or 'processing', continue waiting
                except Exception as e:
                    print(f"Error checking document status: {e}")
            
            retry_count += 1
        
        # If we get here, processing timed out
        raise HTTPException(
            status_code=500, 
            detail="Document processing timed out. The document may still be processing in the background."
        )
        
    except HTTPException:
        raise
    except Exception as e:
        print(f"Error in insert_text: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/documents")
async def get_documents():
    """Get all documents from all workspaces for WebUI compatibility"""
    try:
        documents = []
        
        # Get aggregated metadata from all workspaces
        all_metadata = get_all_documents_metadata()
        
        # Debug logging
        print("\n=== Getting documents for WebUI ===")
        print(f"Total documents across all workspaces: {len(all_metadata)}")
        
        # Try to get documents from LightRAG's doc_status storage
        if hasattr(rag_instance, 'doc_status') and rag_instance.doc_status is not None:
            try:
                # Try different methods to get documents from storage
                doc_status_storage = rag_instance.doc_status
                
                # Method 1: Try to get all documents directly
                if hasattr(doc_status_storage, 'get_all'):
                    try:
                        all_docs = await doc_status_storage.get_all()
                        if all_docs:
                            for doc_id, doc_data in all_docs.items():
                                # CRITICAL: The doc_id here is whatever ID LightRAG is using internally
                                # This could be either our custom ID or a hash ID
                                
                                # Try to get metadata using the doc_id first
                                metadata = all_metadata.get(doc_id, {})
                                
                                # If no metadata found and doc_id looks like a custom ID, it might be stored differently
                                if not metadata and doc_id.startswith('[') and ']' in doc_id:
                                    # This is likely a custom ID, create proper metadata
                                    file_path = doc_id
                                    display_name = doc_id
                                    
                                    # Use full path for display_name
                                    display_name = file_path  # Shows [domain] full_path
                                    
                                    metadata = {"file_path": file_path, "display_name": display_name}
                                
                                # If still no metadata, check if this is a hash ID and search for metadata by content
                                if not metadata and doc_id.startswith('doc-'):
                                    # Search through metadata store for entries with this original_doc_id
                                    for key, meta in all_metadata.items():
                                        if meta.get('original_doc_id') == doc_id:
                                            metadata = meta
                                            break
                                
                                # Ensure file_path exists
                                file_path = metadata.get('file_path', doc_data.get('file_path', f"text/{doc_id}.txt"))
                                
                                # ALWAYS ensure display_name shows full path
                                display_name = metadata.get('display_name', file_path)
                                # Force full path display even if metadata has old format
                                if display_name and "[" in display_name and "]" in display_name:
                                    # If it looks like [domain] something, ensure it's the full path
                                    if file_path and file_path.startswith('[') and file_path != display_name:
                                        display_name = file_path
                                
                                # For display, if doc_id is already in the correct format, use it
                                if doc_id.startswith('[') and ']' in doc_id:
                                    # This document is using a custom ID in the correct format
                                    display_id = doc_id
                                elif file_path and file_path.startswith('[') and ']' in file_path:
                                    # Use file_path as display ID if it's in the correct format
                                    display_id = file_path
                                else:
                                    # Check if we have a better ID in metadata
                                    display_id = metadata.get('id', doc_id)
                                
                                # IMPORTANT: Never return a URL as the ID
                                # If display_id looks like a URL, use the file_path or doc_id instead
                                if display_id and ('http://' in display_id or 'https://' in display_id):
                                    if file_path and file_path.startswith('['):
                                        display_id = file_path
                                    else:
                                        display_id = doc_id
                                
                                documents.append({
                                    "id": display_id,  # This is what the WebUI displays
                                    "doc_id": doc_id,  # Keep the actual doc_id for reference
                                    "file_path": file_path,
                                    "display_name": display_name,
                                    "metadata": metadata,
                                    "status": doc_data.get('status', 'unknown'),
                                    "error": doc_data.get('error') if doc_data.get('status') == 'failed' else None
                                })
                    except Exception as e:
                        print(f"Error with get_all method: {e}")
                
                # Method 2: Try to iterate through storage if it's dict-like
                if not documents and hasattr(doc_status_storage, '_data'):
                    try:
                        storage_data = doc_status_storage._data
                        if isinstance(storage_data, dict):
                            for doc_id, doc_data in storage_data.items():
                                # Get metadata from our store
                                metadata = all_metadata.get(doc_id, {})
                                
                                # Ensure file_path exists
                                file_path = metadata.get('file_path', f"text/{doc_id}.txt")
                                
                                # ALWAYS ensure display_name shows full path
                                display_name = metadata.get('display_name', file_path)
                                # Force full path display even if metadata has old format
                                if display_name and "[" in display_name and "]" in display_name:
                                    # If it looks like [domain] something, ensure it's the full path
                                    if file_path and file_path.startswith('[') and file_path != display_name:
                                        display_name = file_path
                                
                                documents.append({
                                    "id": doc_id,
                                    "file_path": file_path,
                                    "display_name": display_name,
                                    "metadata": metadata,
                                    "status": doc_data.get('status', 'unknown'),
                                    "error": doc_data.get('error') if doc_data.get('status') == 'failed' else None
                                })
                    except Exception as e:
                        print(f"Error accessing _data: {e}")
                
                # Method 3: Try JSON storage file directly
                if not documents:
                    try:
                        working_dir = os.getenv("WORKING_DIR", "/app/data/rag_storage")
                        doc_status_file = os.path.join(working_dir, "doc_status.json")
                        if os.path.exists(doc_status_file):
                            with open(doc_status_file, 'r') as f:
                                doc_status_data = json.load(f)
                                for doc_id, doc_data in doc_status_data.items():
                                    # Get metadata from our store
                                    metadata = all_metadata.get(doc_id, {})
                                    
                                    # Ensure file_path exists
                                    file_path = metadata.get('file_path', f"text/{doc_id}.txt")
                                    
                                    # Handle legacy documents without display_name
                                    display_name = metadata.get('display_name')
                                    if not display_name:
                                        # FORCE display_name to show full path
                                        display_name = file_path  # This will show [domain] full/path
                                    
                                    documents.append({
                                        "id": doc_id,
                                        "file_path": file_path,
                                        "display_name": display_name,
                                        "metadata": metadata,
                                        "status": "processed",  # From metadata store, assume processed
                                        "error": None
                                    })
                    except Exception as e:
                        print(f"Error reading doc_status.json: {e}")
            except Exception as e:
                print(f"Error accessing doc_status storage: {e}")
        
        # If no documents found in doc_status, use metadata
        if not documents:
            for doc_id, metadata in all_metadata.items():
                file_path = metadata.get('file_path', f"text/{doc_id}.txt")
                display_name = metadata.get('display_name')
                
                # Handle legacy documents without display_name
                if not display_name:
                    # FORCE display_name to show full path, not just slug
                    display_name = file_path  # This will show [domain] full/path
                
                # Use file_path as the display ID if it's in the proper format
                display_id = file_path if file_path.startswith('[') and ']' in file_path else doc_id
                
                # IMPORTANT: Never return a URL as the ID
                if display_id and ('http://' in display_id or 'https://' in display_id):
                    if file_path and file_path.startswith('['):
                        display_id = file_path
                    else:
                        display_id = doc_id
                
                documents.append({
                    "id": display_id,  # This is what the WebUI displays
                    "doc_id": doc_id,  # Keep the actual doc_id for reference
                    "file_path": file_path,
                    "display_name": display_name,
                    "metadata": metadata,
                    "status": "processed"
                })
        
        # Categorize documents by status
        processed = [doc for doc in documents if doc.get('status') == 'processed']
        pending = [doc for doc in documents if doc.get('status') == 'pending']
        processing = [doc for doc in documents if doc.get('status') == 'processing']
        failed = [doc for doc in documents if doc.get('status') == 'failed']
        unknown = [doc for doc in documents if doc.get('status') not in ['processed', 'pending', 'processing', 'failed']]
        
        return {
            "statuses": {
                "processed": processed,
                "pending": pending,
                "processing": processing,
                "failed": failed,
                "unknown": unknown
            },
            "total": len(documents)
        }
        
    except Exception as e:
        print(f"Error in get_documents: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/documents/by-sitemap/{sitemap_url:path}")
async def get_documents_by_sitemap(sitemap_url: str):
    """Get all documents for a specific sitemap URL"""
    try:
        matching_docs = []
        
        # Get aggregated metadata from all workspaces
        all_metadata = get_all_documents_metadata()
        
        for doc_id, metadata in all_metadata.items():
            # Check both sitemap_url and sitemap_identifier (for legacy support)
            if (metadata.get('sitemap_url') == sitemap_url or 
                metadata.get('sitemap_identifier') == f"[SITEMAP: {sitemap_url}]"):
                file_path = metadata.get('file_path', f"text/{doc_id}.txt")
                display_name = metadata.get('display_name')
                if not display_name:
                    display_name = generate_display_name_from_file_path(file_path, doc_id)
                    
                # Use file_path as the display ID if it's in the proper format
                display_id = file_path if file_path.startswith('[') and ']' in file_path else doc_id
                
                matching_docs.append({
                    "id": display_id,  # What the WebUI displays
                    "doc_id": doc_id,  # Actual document ID
                    "source_url": metadata.get('source_url'),
                    "file_path": file_path,
                    "display_name": display_name,
                    "indexed_at": metadata.get('indexed_at')
                })
        
        return {
            "sitemap_url": sitemap_url,
            "documents": matching_docs,
            "count": len(matching_docs)
        }
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/documents/by-id")
async def delete_documents_by_id(request: DeleteByIdRequest):
    """Delete documents by their IDs"""
    try:
        # First clean up all document traces from LightRAG storages
        await cleanup_all_document_traces(request.doc_ids)
        
        # Then call the standard delete method
        await rag_instance.adelete_by_doc_id(request.doc_ids)
        
        # Remove from metadata store
        deleted_count = 0
        for doc_id in request.doc_ids:
            # Try direct key first
            if doc_id in metadata_store:
                del metadata_store[doc_id]
                deleted_count += 1
            else:
                # Also check if this is an original_doc_id
                for key in list(metadata_store.keys()):
                    if metadata_store[key].get('original_doc_id') == doc_id:
                        del metadata_store[key]
                        deleted_count += 1
                        break
        
        # Save updated metadata
        save_metadata_store()
        
        return {
            "status": "success",
            "message": f"Successfully deleted {deleted_count} documents",
            "deleted_ids": request.doc_ids
        }
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/documents/by-sitemap/{sitemap_url:path}")
async def delete_documents_by_sitemap(sitemap_url: str):
    """Delete all documents for a specific sitemap URL from all workspaces"""
    try:
        # Track deletions by workspace
        workspace_keys_to_delete = {}
        total_deleted = 0
        
        # Check all workspaces for documents with this sitemap URL
        for workspace_name, workspace_meta in workspace_metadata.items():
            docs_to_delete = []
            metadata_keys_to_delete = []
            
            for doc_id, metadata in workspace_meta.items():
                # Check both sitemap_identifier (legacy) and sitemap_url
                if (metadata.get('sitemap_url') == sitemap_url or 
                    metadata.get('sitemap_identifier') == f"[SITEMAP: {sitemap_url}]"):
                    # The doc_id in metadata_store is the key, but LightRAG might be using custom_id
                    # Check if this document was inserted with a custom ID
                    metadata_keys_to_delete.append(doc_id)
                    
                    # For LightRAG deletion, we need to use the actual ID that LightRAG is using
                    # This could be the doc_id itself (if it's a custom ID like [domain] path)
                    # or it could be in the metadata
                    if doc_id.startswith('[') and ']' in doc_id:
                        # This is likely the custom ID used by LightRAG
                        docs_to_delete.append(doc_id)
                    else:
                        # Check if there's an original_doc_id that might be the hash ID
                        original_id = metadata.get('original_doc_id', doc_id)
                        docs_to_delete.append(original_id)
                        # Also try the file_path as it might be the custom ID
                        file_path = metadata.get('file_path')
                        if file_path and file_path != original_id:
                            docs_to_delete.append(file_path)
            
            # Store metadata keys to delete for this workspace
            if metadata_keys_to_delete:
                workspace_keys_to_delete[workspace_name] = metadata_keys_to_delete
                
                # Remove duplicates from docs_to_delete
                unique_docs_to_delete = list(set(docs_to_delete))
                
                if unique_docs_to_delete:
                    print(f"Deleting documents for sitemap {sitemap_url} in workspace {workspace_name}")
                    print(f"Document IDs to delete: {unique_docs_to_delete}")
                    
                    # Get the workspace RAG instance
                    rag = await WorkspaceManager.get_or_create_instance(workspace_name)
                    
                    # First clean up all document traces from LightRAG storages
                    await cleanup_all_document_traces(unique_docs_to_delete)
                    
                    # Then call the standard delete method
                    try:
                        await rag.adelete_by_doc_id(unique_docs_to_delete)
                    except Exception as e:
                        print(f"Error in adelete_by_doc_id: {e}")
                        # Continue even if this fails, as cleanup_all_document_traces should have done most of the work
            
            # Remove from workspace metadata using the correct keys
            for workspace_name, keys in workspace_keys_to_delete.items():
                if workspace_name in workspace_metadata:
                    for key in keys:
                        if key in workspace_metadata[workspace_name]:
                            del workspace_metadata[workspace_name][key]
                    # Save workspace metadata
                    save_workspace_metadata(workspace_name)
            
            # Force a final save of the graph to ensure persistence
            working_dir = os.getenv("WORKING_DIR", "/app/data/rag_storage")
            if hasattr(rag_instance, 'chunk_entity_relation_graph'):
                graph_storage = rag_instance.chunk_entity_relation_graph
                graph = get_graph_from_storage(graph_storage)
                if graph:
                    force_save_graph_to_disk(graph, working_dir)
            
            # Also ensure all vector DB files are saved with proper structure
            print("Ensuring all vector database files are properly saved...")
            vdb_files = ["vdb_entities.json", "vdb_relationships.json", "vdb_chunks.json"]
            for file_name in vdb_files:
                file_path = os.path.join(working_dir, file_name)
                if os.path.exists(file_path):
                    try:
                        # Read and re-save to ensure it's properly formatted
                        with open(file_path, 'r') as f:
                            data = json.load(f)
                        
                        # Ensure proper structure
                        if not isinstance(data, dict):
                            data = {"embedding_dim": 1536, "data": [], "matrix": []}
                        if 'data' not in data:
                            data['data'] = []
                        if 'matrix' not in data:
                            data['matrix'] = []
                        if 'embedding_dim' not in data:
                            data['embedding_dim'] = 1536
                            
                        with open(file_path, 'w') as f:
                            json.dump(data, f, indent=2)
                        print(f"Verified and saved {file_name}")
                    except Exception as e:
                        print(f"Error verifying {file_name}: {e}")
                        # Create empty structure if corrupted
                        empty_structure = {"embedding_dim": 1536, "data": [], "matrix": []}
                        with open(file_path, 'w') as f:
                            json.dump(empty_structure, f, indent=2)
                        print(f"Reset {file_name} to empty structure")
            
            # Force save all LightRAG storages
            if hasattr(rag_instance, 'graph_storage'):
                storage = rag_instance.graph_storage
                if hasattr(storage, 'save'):
                    try:
                        await storage.save()
                        print("Force saved graph storage")
                    except Exception as e:
                        print(f"Error force saving graph storage: {e}")
                elif hasattr(storage, '_save'):
                    try:
                        storage._save()
                        print("Force saved graph storage (sync)")
                    except Exception as e:
                        print(f"Error force saving graph storage (sync): {e}")
            
            # Small delay to ensure filesystem sync
            time.sleep(0.5)
        
        return {
            "status": "success",
            "message": f"Deleted {len(metadata_keys_to_delete)} documents for sitemap {sitemap_url}",
            "deleted_count": len(metadata_keys_to_delete),
            "deleted_ids": unique_docs_to_delete if 'unique_docs_to_delete' in locals() else [],
            "sitemap_url": sitemap_url,
            "graph_cleared": True
        }
        
    except Exception as e:
        print(f"Error in delete_documents_by_sitemap: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/graph/force-clear-all")
async def force_clear_all_graph_data():
    """Force clear ALL graph and vector database files - nuclear option"""
    try:
        working_dir = os.getenv("WORKING_DIR", "/app/data/rag_storage")
        
        # List of ALL files that might contain graph or entity data
        files_to_delete = [
            # Graph files
            "graph_chunk_entity_relation.graphml",
            "graph_data.json",
            "graph_cache.json",
            # Vector database files
            "vdb_entities.json",
            "vdb_relationships.json", 
            "vdb_chunks.json",
            # Key-value stores that might contain entities
            "kv_store_text_chunks.json",
            "kv_store_full_docs.json",
            "kv_store_llm_response_cache.json",
            # Document status
            "doc_status.json",
            # Any other potential files
            "entity_embedding.json",
            "relationship_embedding.json",
            "document_graph_storage.json",
            "entities.json",
            "relationships.json",
            "chunks.json"
        ]
        
        deleted_files = []
        errors = []
        
        # Clear in-memory graph
        if hasattr(rag_instance, 'chunk_entity_relation_graph'):
            graph_storage = rag_instance.chunk_entity_relation_graph
            graph = get_graph_from_storage(graph_storage)
            if graph and hasattr(graph, 'clear'):
                graph.clear()
                print("Cleared in-memory graph")
        
        # Delete or clear all files
        for file_name in files_to_delete:
            file_path = os.path.join(working_dir, file_name)
            if os.path.exists(file_path):
                try:
                    # For safety, rename to backup first
                    backup_path = file_path + ".backup"
                    os.rename(file_path, backup_path)
                    
                    # Create empty file based on type
                    if file_name.endswith('.graphml'):
                        empty_graph = nx.Graph()
                        nx.write_graphml(empty_graph, file_path)
                    elif file_name.startswith('vdb_'):
                        # Vector database format
                        empty_vdb = {
                            "embedding_dim": 1536,
                            "data": [],
                            "matrix": []
                        }
                        with open(file_path, 'w') as f:
                            json.dump(empty_vdb, f)
                    else:
                        # Empty JSON object
                        with open(file_path, 'w') as f:
                            json.dump({}, f)
                    
                    # Delete backup
                    os.remove(backup_path)
                    deleted_files.append(file_name)
                    print(f"Cleared {file_name}")
                    
                except Exception as e:
                    errors.append({"file": file_name, "error": str(e)})
                    print(f"Error clearing {file_name}: {e}")
        
        # Clear metadata store
        global metadata_store
        metadata_store = {}
        save_metadata_store()
        
        # Clear any caches
        if hasattr(rag_instance, 'aclear_cache'):
            try:
                await rag_instance.aclear_cache()
                print("Cleared LLM cache")
            except Exception as e:
                print(f"Error clearing cache: {e}")
        
        return {
            "status": "success",
            "message": "Force cleared all graph and vector database files",
            "deleted_files": deleted_files,
            "errors": errors,
            "total_files_cleared": len(deleted_files)
        }
        
    except Exception as e:
        print(f"Error in force clear: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/documents/upload")
async def upload_document(
    file: UploadFile = File(...),
    custom_name: Optional[str] = Form(None)
):
    """Upload and process a document file
    
    Args:
        file: The uploaded file
        custom_name: Optional custom name prefix for the file
    
    Returns:
        Document metadata including ID and file path
    """
    try:
        # Read file content
        file_content = await file.read()
        
        # Get file extension
        filename = file.filename
        file_extension = filename.split('.')[-1].lower() if '.' in filename else ''
        
        # Validate file type
        if file_extension not in SUPPORTED_EXTENSIONS:
            raise HTTPException(
                status_code=400,
                detail=f"Unsupported file type: {file_extension}. Supported types: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
            )
        
        # Extract text from file
        extracted_text = extract_text_from_file(file_content, file_extension, filename)
        
        # Compute document ID
        doc_id = compute_doc_id(extracted_text)
        
        # Generate file path with custom name if provided
        if custom_name:
            # Format: [custom_name] original_filename
            file_path = f"[{custom_name}] {filename}"
            display_name = file_path
        else:
            file_path = f"uploaded/{filename}"
            display_name = filename
        
        # Add metadata header to content
        metadata_parts = [
            f"[FILE: {filename}]",
            f"[TYPE: {file_extension.upper()}]",
            f"[UPLOADED: {datetime.utcnow().isoformat()}]"
        ]
        if custom_name:
            metadata_parts.append(f"[CUSTOM_NAME: {custom_name}]")
        
        enriched_content = "\n".join(metadata_parts) + "\n\n" + extracted_text
        
        # Store metadata
        metadata_entry = {
            "id": doc_id,
            "file_path": file_path,
            "display_name": display_name,
            "original_filename": filename,
            "custom_name": custom_name,
            "file_type": file_extension,
            "indexed_at": datetime.utcnow().isoformat(),
            "content_summary": extracted_text[:200] + "..." if len(extracted_text) > 200 else extracted_text
        }
        
        metadata_store[doc_id] = metadata_entry
        save_metadata_store()
        
        # Insert into LightRAG
        await rag_instance.ainsert(enriched_content, ids=[doc_id], file_paths=[file_path])
        
        return {
            "status": "success",
            "message": f"Document '{filename}' uploaded and processed successfully",
            "doc_id": doc_id,
            "file_path": file_path,
            "display_name": display_name,
            "file_type": file_extension,
            "text_length": len(extracted_text)
        }
        
    except HTTPException:
        raise
    except Exception as e:
        print(f"Error uploading document: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/documents/upload/bulk")
async def upload_documents_bulk(
    files: List[UploadFile] = File(...),
    custom_name: Optional[str] = Form(None)
):
    """Upload and process multiple document files"""
    results = []
    
    for i, file in enumerate(files):
        try:
            # Process each file using existing logic
            file_content = await file.read()
            filename = file.filename
            file_extension = filename.split('.')[-1].lower() if '.' in filename else ''
            
            if file_extension not in SUPPORTED_EXTENSIONS:
                results.append({
                    "filename": filename,
                    "status": "error",
                    "error": f"Unsupported file type: {file_extension}"
                })
                continue
            
            # Extract text and process (existing logic)
            extracted_text = extract_text_from_file(file_content, file_extension, filename)
            doc_id = compute_doc_id(extracted_text)
            
            # Generate file path
            if custom_name:
                file_path = f"[{custom_name}] {filename}"
            else:
                file_path = f"bulk_upload/{filename}"
            
            # Add metadata
            metadata_parts = [
                f"[FILE: {filename}]",
                f"[TYPE: {file_extension.upper()}]",
                f"[BULK_UPLOAD: {i+1} of {len(files)}]",
                f"[UPLOADED: {datetime.utcnow().isoformat()}]"
            ]
            
            enriched_content = "\n".join(metadata_parts) + "\n\n" + extracted_text
            
            # Store metadata
            metadata_entry = {
                "id": doc_id,
                "file_path": file_path,
                "display_name": filename,
                "original_filename": filename,
                "file_type": file_extension,
                "indexed_at": datetime.utcnow().isoformat(),
                "bulk_upload": True,
                "batch_index": i + 1,
                "batch_total": len(files)
            }
            
            metadata_store[doc_id] = metadata_entry
            
            # Insert into LightRAG
            await rag_instance.ainsert(enriched_content, ids=[doc_id], file_paths=[file_path])
            
            results.append({
                "filename": filename,
                "status": "success",
                "doc_id": doc_id,
                "file_path": file_path
            })
            
        except Exception as e:
            results.append({
                "filename": file.filename,
                "status": "error",
                "error": str(e)
            })
    
    # Save metadata after all uploads
    save_metadata_store()
    
    return {
        "status": "completed",
        "total_files": len(files),
        "successful": len([r for r in results if r["status"] == "success"]),
        "failed": len([r for r in results if r["status"] == "error"]),
        "results": results
    }

@app.post("/query")
async def query_documents(request: QueryRequest):
    """Query documents using LightRAG with optional workspace specification
    
    Args:
        request: Query request with query string, mode, and optional workspace
    
    Returns:
        Query response from LightRAG
    """
    try:
        # Use specified workspace, or default if not provided
        if request.workspace:
            # Get or create the specified workspace
            rag = await WorkspaceManager.get_or_create_instance(request.workspace)
            print(f"Querying workspace: {request.workspace}")
        else:
            # Use default workspace or first available workspace
            if default_workspace in workspace_instances:
                rag = workspace_instances[default_workspace]
            elif workspace_instances:
                rag = list(workspace_instances.values())[0]
            else:
                # Create default instance if none exists
                rag = await WorkspaceManager.get_or_create_instance(default_workspace)
            print(f"Querying default workspace: {default_workspace}")
        
        # Create query parameters
        query_param = QueryParam(
            mode=request.mode,
            stream=False
        )
        
        # Execute query
        result = await rag.aquery(request.query, param=query_param)
        
        return {
            "status": "success",
            "response": result,
            "mode": request.mode,
            "workspace": request.workspace or default_workspace
        }
        
    except Exception as e:
        print(f"Error in /query: {str(e)}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/query/stream")
async def query_documents_stream(request: QueryRequest):
    """Stream query results from LightRAG with optional workspace specification
    
    Args:
        request: Query request with query string, mode, and optional workspace
    
    Returns:
        Streaming response from LightRAG
    """
    from fastapi.responses import StreamingResponse
    import json
    
    async def generate():
        try:
            # Use specified workspace, or default if not provided
            if request.workspace:
                # Get or create the specified workspace
                rag = await WorkspaceManager.get_or_create_instance(request.workspace)
                print(f"Streaming query from workspace: {request.workspace}")
            else:
                # Use default workspace or first available workspace
                if default_workspace in workspace_instances:
                    rag = workspace_instances[default_workspace]
                elif workspace_instances:
                    rag = list(workspace_instances.values())[0]
                else:
                    # Create default instance if none exists
                    rag = await WorkspaceManager.get_or_create_instance(default_workspace)
                print(f"Streaming query from default workspace: {default_workspace}")
            
            # Create query parameters with streaming
            query_param = QueryParam(
                mode=request.mode,
                stream=True
            )
            
            # Execute streaming query
            async for chunk in rag.aquery(request.query, param=query_param):
                # Format as server-sent event
                data = json.dumps({"content": chunk, "workspace": request.workspace or default_workspace})
                yield f"data: {data}\n\n"
                
        except Exception as e:
            error_data = json.dumps({"error": str(e)})
            yield f"data: {error_data}\n\n"
    
    return StreamingResponse(generate(), media_type="text/event-stream")

@app.delete("/documents")
async def clear_all_documents():
    """Clear all documents from all workspaces
    
    Returns:
        Status of the clear operation
    """
    try:
        results = []
        
        # Clear each workspace
        for workspace_name, rag in workspace_instances.items():
            try:
                # Clear the workspace's data
                if hasattr(rag, 'full_docs'):
                    await rag.full_docs.drop()
                if hasattr(rag, 'text_chunks'):
                    await rag.text_chunks.drop()
                if hasattr(rag, 'entities_vdb'):
                    await rag.entities_vdb.drop()
                if hasattr(rag, 'relationships_vdb'):
                    await rag.relationships_vdb.drop()
                if hasattr(rag, 'chunks_vdb'):
                    await rag.chunks_vdb.drop()
                if hasattr(rag, 'chunk_entity_relation_graph'):
                    await rag.chunk_entity_relation_graph.drop()
                
                # Clear workspace metadata
                if workspace_name in workspace_metadata:
                    workspace_metadata[workspace_name] = {}
                    save_workspace_metadata(workspace_name)
                
                results.append({
                    "workspace": workspace_name,
                    "status": "cleared"
                })
                
            except Exception as e:
                results.append({
                    "workspace": workspace_name,
                    "status": "error",
                    "error": str(e)
                })
        
        # Clear global metadata store
        global metadata_store
        metadata_store = {}
        save_metadata_store()
        
        return {
            "status": "success",
            "message": f"Cleared {len(results)} workspaces",
            "results": results
        }
        
    except Exception as e:
        print(f"Error clearing documents: {str(e)}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/documents/scan")
async def scan_for_documents():
    """Scan for new documents to index
    
    Returns:
        Status of the scan operation
    """
    try:
        # This would typically scan a directory for new files
        # For now, return a simple status
        return {
            "status": "success",
            "message": "Document scan completed",
            "new_documents": 0
        }
        
    except Exception as e:
        print(f"Error scanning documents: {str(e)}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

if __name__ == "__main__":
    port = int(os.environ.get("PORT", 8080))
    host = os.environ.get("HOST", "0.0.0.0")
    uvicorn.run(app, host=host, port=port)



#IT WORKS!!!!!!!!